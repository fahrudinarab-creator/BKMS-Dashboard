import streamlit as st
import pandas as pd
import plotly.graph_objects as go
import plotly.express as px
import plotly.io as pio
import base64
import re
from pathlib import Path

pio.templates.default = "plotly_dark"

_unit_code_pattern = re.compile(r'(\d{3}-\d{3})\s*$')
def _unit_label(name):
    if not isinstance(name, str):
        return name
    m = _unit_code_pattern.search(name)
    code = m.group(1) if m else "?"
    return f"{code} — {name}"

def style_fig(fig):
    """Make chart background transparent so it blends with the black page,
    and keep text light/readable."""
    fig.update_layout(
        paper_bgcolor="rgba(0,0,0,0)",
        plot_bgcolor="rgba(0,0,0,0)",
        font=dict(color="#F3F4F6"),
        legend=dict(font=dict(color="#F3F4F6")),
    )
    fig.update_xaxes(gridcolor="#2D333B", zerolinecolor="#2D333B")
    fig.update_yaxes(gridcolor="#2D333B", zerolinecolor="#2D333B")
    return fig

# ---------------------------------------------------------------
# PAGE CONFIG
# ---------------------------------------------------------------
st.set_page_config(
    page_title="Dashboard Operational Review | PT BKMS",
    page_icon="📊",
    layout="wide",
    initial_sidebar_state="expanded",
)

# ---------------------------------------------------------------
# STYLE
# ---------------------------------------------------------------
PRIMARY = "#0B3D2E"      # forest green (used for banner/backgrounds)
CHART_GREEN = "#3FA772"  # brighter green for readable bars/lines on black
GOLD = "#C9A227"         # gold accent
RED = "#E4574C"          # brighter red for readability on black
GREY = "#9CA3AF"         # lighter grey for readability on black
EMOJI_FONT = "Segoe UI Emoji"  # font khusus utk karakter emoji, supaya render benar di PowerPoint asli
                                # (Calibri/font teks biasa tidak punya glyph emoji -> muncul kotak/silang)

DARK_BG = "#000000"
CARD_BG = "#161B22"
BORDER = "#2D333B"
TEXT_LIGHT = "#F3F4F6"
TEXT_MUTED = "#9CA3AF"

st.markdown(f"""
<style>
    html, body, [data-testid="stAppViewContainer"], .main {{
        background-color: {DARK_BG} !important;
    }}
    [data-testid="stHeader"] {{ background-color: rgba(0,0,0,0) !important; }}
    [data-testid="stSidebar"] {{ background-color: {CARD_BG} !important; border-right: 1px solid {BORDER}; }}
    [data-testid="stSidebar"] * {{ color: {TEXT_LIGHT} !important; }}
    .block-container {{ padding-top: 1.5rem; }}

    div[data-testid="stMetric"] {{
        background: {CARD_BG} !important;
        border: 1px solid {BORDER};
        border-left: 5px solid {GOLD};
        border-radius: 10px;
        padding: 14px 16px 10px 16px;
    }}
    div[data-testid="stMetricLabel"] * {{ color: {TEXT_MUTED} !important; font-weight: 600; }}
    div[data-testid="stMetricValue"] * {{ color: {TEXT_LIGHT} !important; font-weight: 700; }}
    div[data-testid="stMetricDelta"] * {{ font-weight: 600; }}

    h1, h2, h3, h4, h5, h6, p, span, label, .stMarkdown {{ color: {TEXT_LIGHT} !important; }}
    h1, h2, h3 {{ color: {GOLD} !important; }}

    .header-banner {{
        background: linear-gradient(90deg, {PRIMARY} 0%, #145C43 100%);
        padding: 22px 200px 22px 28px;
        border-radius: 12px;
        margin-bottom: 18px;
        border: 1px solid {BORDER};
        position: relative;
        min-height: 160px;
        display: flex;
        flex-direction: column;
        justify-content: center;
    }}
    .header-banner h1 {{ color: white !important; margin: 0; font-size: 26px; }}
    .header-banner p {{ color: {GOLD} !important; margin: 2px 0 0 0; font-size: 14px; letter-spacing: 0.5px; }}
    .header-logo {{
        position: absolute;
        top: 50%;
        right: 28px;
        transform: translateY(-50%);
        height: 170px;
        width: auto;
    }}
    .section-title {{
        padding-left: 4px;
        margin-top: 6px;
        color: {TEXT_LIGHT} !important;
    }}
    .insight-box {{
        background: {CARD_BG};
        border: 1px solid {BORDER};
        border-left: 5px solid {CHART_GREEN};
        border-radius: 10px;
        padding: 16px 20px;
        margin-bottom: 10px;
    }}
    .insight-box li {{ margin-bottom: 8px; line-height: 1.5; }}

    /* KPI cards (icon + big number + status pill), styled after the RTM report format */
    .kpi-card {{
        background: {CARD_BG};
        border: 1px solid {BORDER};
        border-top: 4px solid var(--accent, {GOLD});
        border-radius: 12px;
        padding: 18px 18px 16px 18px;
        height: 100%;
    }}
    .kpi-icon {{
        width: 40px; height: 40px; border-radius: 50%;
        display: flex; align-items: center; justify-content: center;
        font-size: 19px; margin-bottom: 10px;
    }}
    .kpi-label {{ color: {TEXT_MUTED} !important; font-size: 13px; font-weight: 600; margin-bottom: 4px; }}
    .kpi-value {{ color: {TEXT_LIGHT} !important; font-size: 28px; font-weight: 800; line-height: 1.15; margin-bottom: 6px; }}
    .kpi-budget {{ color: {TEXT_MUTED} !important; font-size: 12.5px; margin-bottom: 10px; }}
    .kpi-pill {{
        display: inline-block; padding: 5px 12px; border-radius: 20px;
        font-size: 12.5px; font-weight: 700;
    }}
    .kpi-pill-green {{ background: rgba(63,167,114,0.18); color: #6EE7A8 !important; }}
    .kpi-pill-red {{ background: rgba(228,87,76,0.18); color: #FF9A91 !important; }}
    .kpi-pill-amber {{ background: rgba(201,162,39,0.18); color: {GOLD} !important; }}

    .ringkasan-table {{
        width: 100%; border-collapse: collapse; border-radius: 10px; overflow: hidden;
        border: 1px solid {BORDER};
    }}
    .ringkasan-table th {{
        background: {PRIMARY}; color: white !important; text-align: left;
        padding: 10px 14px; font-size: 13px;
    }}
    .ringkasan-table td {{
        padding: 10px 14px; font-size: 13.5px; color: {TEXT_LIGHT} !important;
        border-top: 1px solid {BORDER}; background: {CARD_BG};
    }}
    .ringkasan-table tr:nth-child(even) td {{ background: #1B212B; }}

    .ringkasan-table-sm {{
        width: 100%; border-collapse: separate; border-spacing: 0; border-radius: 8px;
        border: 1px solid {BORDER};
    }}
    .ringkasan-table-sm th {{
        background: {PRIMARY}; color: white !important; text-align: left;
        padding: 8px 12px; font-size: 12.5px; white-space: nowrap;
        position: sticky; top: 0; z-index: 2;
    }}
    .ringkasan-table-sm td {{
        padding: 8px 12px; font-size: 13px; color: {TEXT_LIGHT} !important;
        border-top: 1px solid {BORDER}; background: {CARD_BG}; white-space: nowrap;
    }}
    .ringkasan-table-sm tr:nth-child(even) td {{ background: #1B212B; }}
    .rk-badge {{
        display: inline-block; padding: 4px 10px; border-radius: 14px;
        font-size: 12.5px; font-weight: 700; text-align: center; min-width: 55px;
    }}
    .rk-green {{ background: rgba(63,167,114,0.20); color: #6EE7A8 !important; }}
    .rk-red {{ background: rgba(228,87,76,0.20); color: #FF9A91 !important; }}
    .rk-badge-sm {{
        display: inline-block; padding: 3px 10px; border-radius: 12px;
        font-size: 12px; font-weight: 700; text-align: center; min-width: 48px;
    }}
    .rk-grey {{ background: rgba(156,163,175,0.20); color: {TEXT_MUTED} !important; }}

    [data-testid="stDataFrame"] {{ background-color: {CARD_BG} !important; }}
    .stTextInput input {{ background-color: {CARD_BG} !important; color: {TEXT_LIGHT} !important; }}
</style>
""", unsafe_allow_html=True)

def kpi_card(icon, icon_bg, accent, label, value, budget_text, pill_text, pill_style):
    """Render one KPI card matching the RTM-report visual style (icon circle, big number, status pill)."""
    return f"""
    <div class="kpi-card" style="--accent:{accent}">
        <div class="kpi-icon" style="background:{icon_bg}">{icon}</div>
        <div class="kpi-label">{label}</div>
        <div class="kpi-value">{value}</div>
        <div class="kpi-budget">{budget_text}</div>
        <span class="kpi-pill {pill_style}">{pill_text}</span>
    </div>
    """

def achievement_pill(pct, higher_is_better=True, target_label="Target"):
    """Return (pill_text, pill_style) for a metric vs its target/budget."""
    if pct is None:
        return (f"{target_label} = 0", "kpi-pill-amber")
    if higher_is_better:
        if pct >= 100:
            return (f"✓ {pct:.1f}% — Tercapai", "kpi-pill-green")
        return (f"✗ {pct:.1f}% vs {target_label}", "kpi-pill-red")
    else:
        if pct <= 100:
            return (f"✓ {pct:.1f}% — Under Budget", "kpi-pill-green")
        return (f"✗ {pct:.1f}% — Over Budget", "kpi-pill-red")

# ---------------------------------------------------------------
# DATA LOADING
# ---------------------------------------------------------------
DATA_PATH = Path(__file__).parent / "data_bkms.csv"
MAINT_DATA_PATH = Path(__file__).parent / "data_maintenance.csv"
SPAREPART_DATA_PATH = Path(__file__).parent / "data_sparepart.csv"
SASARAN_MUTU_PATH = Path(__file__).parent / "data_sasaran_mutu.csv"
MONTH_ORDER = ["Jan", "Feb", "Mar", "Apr", "May", "Jun", "Jul", "Aug", "Sep", "Oct", "Nov", "Dec"]
KATEGORI_LABEL = {"AB": "Alat Berat (AB)", "TR": "Transportasi (TR)"}

@st.cache_data
def load_data(file) -> pd.DataFrame:
    return pd.read_csv(file)

@st.cache_data
def load_maintenance_data(file) -> pd.DataFrame:
    if not Path(file).exists():
        return pd.DataFrame()
    return pd.read_csv(file)

@st.cache_data
def load_sparepart_data(file) -> pd.DataFrame:
    if not Path(file).exists():
        return pd.DataFrame()
    return pd.read_csv(file)

@st.cache_data
def load_sasaran_mutu_data(file) -> pd.DataFrame:
    if not Path(file).exists():
        return pd.DataFrame()
    return pd.read_csv(file, dtype={"id_unit": str})

def load_from_upload(uploaded_file) -> pd.DataFrame:
    """Parse an uploaded 'Gabungan.xlsx' file with the same fixed layout used to build data_bkms.csv.
    Layout (1-indexed columns): A=ID Unit, B=Kode Unit, C=Nama Unit, D=Nilai Asset,
    E=Kriteria Unit, F=Jenis Unit, G-H=Prestasi (R/B), I-J=Pendapatan (R/B), K-L=Upah (R/B),
    M-N=Qty BBM (R/B), O-P=Harga BBM (R/B), Q-R=Biaya BBM (R/B), S-T=Maintenance (R/B),
    U-V=Penyusutan (R/B), W-X=Lainnya (R/B), Y-Z=Biaya Langsung (R/B), AA-AB=Biaya Tdk Langsung (R/B),
    AC-AD=Total Biaya (R/B), AE=Lokasi, AF=Status, AG=Kategori."""
    import openpyxl
    wb = openpyxl.load_workbook(uploaded_file, data_only=True)
    ws = wb[wb.sheetnames[0]]
    month_map = {'Jan': 1, 'Feb': 2, 'Mar': 3, 'Apr': 4, 'May': 5, 'Jun': 6,
                 'Jul': 7, 'Aug': 8, 'Sep': 9, 'Oct': 10, 'Nov': 11, 'Dec': 12}
    rows = []
    for r in range(5, ws.max_row + 1):
        id_unit = ws.cell(row=r, column=1).value
        if not id_unit:
            continue
        lokasi = ws.cell(row=r, column=31).value
        bulan_nama = ws.cell(row=r, column=32).value
        kategori = ws.cell(row=r, column=33).value
        if not lokasi or not bulan_nama:
            continue
        rows.append(dict(
            id_unit=id_unit, kode_unit=ws.cell(row=r, column=2).value,
            nama_unit=ws.cell(row=r, column=3).value, nilai_asset=ws.cell(row=r, column=4).value or 0,
            kriteria_unit=ws.cell(row=r, column=5).value, jenis_unit=ws.cell(row=r, column=6).value,
            lokasi=lokasi, bulan=bulan_nama, bulan_no=month_map.get(bulan_nama, 0), kategori=kategori,
            prestasi_realisasi=ws.cell(row=r, column=7).value or 0, prestasi_budget=ws.cell(row=r, column=8).value or 0,
            pendapatan_realisasi=ws.cell(row=r, column=9).value or 0, pendapatan_budget=ws.cell(row=r, column=10).value or 0,
            upah_realisasi=ws.cell(row=r, column=11).value or 0, upah_budget=ws.cell(row=r, column=12).value or 0,
            qty_bbm_realisasi=ws.cell(row=r, column=13).value or 0, qty_bbm_budget=ws.cell(row=r, column=14).value or 0,
            harga_bbm_realisasi=ws.cell(row=r, column=15).value or 0, harga_bbm_budget=ws.cell(row=r, column=16).value or 0,
            biaya_bbm_realisasi=ws.cell(row=r, column=17).value or 0, biaya_bbm_budget=ws.cell(row=r, column=18).value or 0,
            maintenance_realisasi=ws.cell(row=r, column=19).value or 0, maintenance_budget=ws.cell(row=r, column=20).value or 0,
            penyusutan_realisasi=ws.cell(row=r, column=21).value or 0, penyusutan_budget=ws.cell(row=r, column=22).value or 0,
            lainnya_realisasi=ws.cell(row=r, column=23).value or 0, lainnya_budget=ws.cell(row=r, column=24).value or 0,
            biaya_langsung_realisasi=ws.cell(row=r, column=25).value or 0, biaya_langsung_budget=ws.cell(row=r, column=26).value or 0,
            biaya_tidak_langsung_realisasi=ws.cell(row=r, column=27).value or 0, biaya_tidak_langsung_budget=ws.cell(row=r, column=28).value or 0,
            total_biaya_realisasi=ws.cell(row=r, column=29).value or 0, total_biaya_budget=ws.cell(row=r, column=30).value or 0,
        ))
    return pd.DataFrame(rows)

def load_from_upload_maintenance(uploaded_file) -> pd.DataFrame:
    """Parse an uploaded maintenance detail file (e.g. 'Pemeliharaan_sd_Bulan.xls').
    DESCRIPTION format: 'PEMELIHARAAN (RUTIN|NON RUTIN) (kategori) (tipe biaya) (PLANTATION|MINING) (site) - (unit)'
    """
    raw = pd.read_excel(uploaded_file, sheet_name=0, header=0)
    pattern = re.compile(
        r'^PEMELIHARAAN\s+(RUTIN|NON RUTIN)\s+(.*?)\s+(SPAREPART|ALOKASI WORKSHOP|SERVICE LUAR|LAIN-LAIN)\s+(PLANTATION|MINING)\s+(.*?)\s+-\s+(.*)$'
    )
    month_id_map = {1: 'Jan', 2: 'Feb', 3: 'Mar', 4: 'Apr', 5: 'May', 6: 'Jun',
                     7: 'Jul', 8: 'Aug', 9: 'Sep', 10: 'Oct', 11: 'Nov', 12: 'Dec'}
    raw["FDATE"] = pd.to_datetime(raw["FDATE"], errors="coerce")
    raw = raw.dropna(subset=["FDATE", "DESCRIPTION", "FAMOUNT_RP"])
    raw = raw[raw["FDATE"].dt.year >= 2000]

    rows = []
    for _, r in raw.iterrows():
        desc = str(r["DESCRIPTION"]).strip()
        m = pattern.match(desc)
        if not m:
            continue
        jenis, kategori_sparepart, tipe_biaya, kelompok, lokasi, unit = m.groups()
        dt = r["FDATE"]
        rows.append(dict(
            tanggal=dt.date().isoformat(), bulan=month_id_map.get(dt.month, ''), bulan_no=dt.month,
            lokasi=lokasi.strip(), kelompok=kelompok.strip(),
            jenis_pemeliharaan=jenis.strip(), kategori_sparepart=kategori_sparepart.strip(),
            tipe_biaya=tipe_biaya.strip(), nama_unit=unit.strip(),
            biaya=float(r["FAMOUNT_RP"]),
            keterangan=str(r["FREMARK"]) if pd.notna(r["FREMARK"]) else '',
        ))
    return pd.DataFrame(rows)

def load_from_upload_sparepart(uploaded_file) -> pd.DataFrame:
    """Parse an uploaded spare-part usage detail file (e.g. 'Rincian_Pemakaian.xls').
    Note: in the source export, the long description text is actually stored in the
    'ACCOUNT_DESC' column (the 'KETERANGAN' column only holds the unit code)."""
    raw = pd.read_excel(uploaded_file, sheet_name=0, header=0)
    pattern = re.compile(
        r'^PEMELIHARAAN\s+(RUTIN|NON RUTIN)\s+(.*?)\s+(SPAREPART|ALOKASI WORKSHOP|SERVICE LUAR|LAIN-LAIN)\s+(PLANTATION|MINING)\s+(.*?)\s+-\s+(.*)$'
    )
    month_id_map = {1: 'Jan', 2: 'Feb', 3: 'Mar', 4: 'Apr', 5: 'May', 6: 'Jun',
                     7: 'Jul', 8: 'Aug', 9: 'Sep', 10: 'Oct', 11: 'Nov', 12: 'Dec'}
    raw["TGL"] = pd.to_datetime(raw["TGL"], errors="coerce")
    raw = raw.dropna(subset=["TGL", "ACCOUNT_DESC", "TOTAL"])

    rows = []
    for _, r in raw.iterrows():
        desc = str(r["ACCOUNT_DESC"]).strip()
        m = pattern.match(desc)
        if not m:
            continue
        jenis, kategori_sparepart, tipe_biaya, kelompok, lokasi, unit = m.groups()
        dt = r["TGL"]
        rows.append(dict(
            tanggal=dt.date().isoformat(), bulan=month_id_map.get(dt.month, ''), bulan_no=dt.month,
            lokasi=lokasi.strip(), kelompok=kelompok.strip(),
            jenis_pemeliharaan=jenis.strip(), kategori_sparepart=kategori_sparepart.strip(),
            nama_unit=unit.strip(), kode_barang=str(r["KODE BRG"]),
            part_number=str(r["PART NUMBER"]) if pd.notna(r["PART NUMBER"]) else '',
            nama_barang=str(r["NAMA BARANG"]), qty=float(r["QTY"]), satuan=str(r["SAT"]),
            biaya=float(r["TOTAL"]),
            group_desc=str(r["GROUP_DESC"]), class_desc=str(r["CLASS_DESC"]), subclass_desc=str(r["SUBCLASS_DESC"]),
        ))
    return pd.DataFrame(rows)

with st.sidebar:
    MINING_SITES = ["TANJUNG", "BUHUT", "BUHUT LHL"]
    PLANTATION_SITES = ["SUNGAI DANAU", "KUMAI"]
    DIVISI_MAP = {"Mining": MINING_SITES, "Plantation": PLANTATION_SITES}

    st.markdown("### 📁 Sumber Data")
    uploaded = st.file_uploader("Upload file Gabungan.xlsx terbaru (opsional)", type=["xlsx"])
    if uploaded is not None:
        try:
            df_raw = load_from_upload(uploaded)
            st.success(f"Berhasil memuat {len(df_raw):,} baris dari file upload.")
        except Exception as e:
            st.error(f"Gagal membaca file: {e}")
            df_raw = load_data(DATA_PATH)
    else:
        df_raw = load_data(DATA_PATH)

    uploaded_maint = st.file_uploader("Upload data Maintenance (Pemeliharaan) terbaru (opsional)", type=["xls", "xlsx"])
    if uploaded_maint is not None:
        try:
            maint_raw = load_from_upload_maintenance(uploaded_maint)
            st.success(f"Berhasil memuat {len(maint_raw):,} baris data maintenance dari file upload.")
        except Exception as e:
            st.error(f"Gagal membaca file maintenance: {e}")
            maint_raw = load_maintenance_data(MAINT_DATA_PATH)
    else:
        maint_raw = load_maintenance_data(MAINT_DATA_PATH)

    uploaded_sparepart = st.file_uploader("Upload data Rincian Pemakaian Sparepart terbaru (opsional)", type=["xls", "xlsx"])
    if uploaded_sparepart is not None:
        try:
            sparepart_raw = load_from_upload_sparepart(uploaded_sparepart)
            st.success(f"Berhasil memuat {len(sparepart_raw):,} baris data pemakaian sparepart dari file upload.")
        except Exception as e:
            st.error(f"Gagal membaca file sparepart: {e}")
            sparepart_raw = load_sparepart_data(SPAREPART_DATA_PATH)
    else:
        sparepart_raw = load_sparepart_data(SPAREPART_DATA_PATH)

    # Tambahkan kolom 'kategori' (AB/TR) ke data maintenance & sparepart, dicocokkan lewat nama_unit
    # terhadap data utama (df_raw) — supaya bisa di-crosscheck per kategori. Hasilnya disimpan kembali
    # ke file CSV-nya (data_maintenance.csv & data_sparepart.csv) supaya kolom kategori permanen di file.
    if not df_raw.empty and "nama_unit" in df_raw.columns and "kategori" in df_raw.columns:
        _kategori_lookup = (
            df_raw.dropna(subset=["nama_unit", "kategori"])
            .assign(_nama_unit_key=lambda d: d["nama_unit"].astype(str).str.strip().str.upper())
            .drop_duplicates(subset=["_nama_unit_key"])
            .set_index("_nama_unit_key")["kategori"]
        )
        if not maint_raw.empty and "nama_unit" in maint_raw.columns:
            maint_raw = maint_raw.copy()
            maint_raw["kategori"] = maint_raw["nama_unit"].astype(str).str.strip().str.upper().map(_kategori_lookup)
            try:
                maint_raw.to_csv(MAINT_DATA_PATH, index=False)
                load_maintenance_data.clear()
            except Exception as _e_maint_save:
                st.warning(f"Kolom kategori berhasil ditambahkan, tapi gagal menyimpan ke {MAINT_DATA_PATH.name}: {_e_maint_save}")
        if not sparepart_raw.empty and "nama_unit" in sparepart_raw.columns:
            sparepart_raw = sparepart_raw.copy()
            sparepart_raw["kategori"] = sparepart_raw["nama_unit"].astype(str).str.strip().str.upper().map(_kategori_lookup)
            try:
                sparepart_raw.to_csv(SPAREPART_DATA_PATH, index=False)
                load_sparepart_data.clear()
            except Exception as _e_sp_save:
                st.warning(f"Kolom kategori berhasil ditambahkan, tapi gagal menyimpan ke {SPAREPART_DATA_PATH.name}: {_e_sp_save}")

    sasaran_mutu_raw = load_sasaran_mutu_data(SASARAN_MUTU_PATH)

    st.markdown("---")
    _download_maint_slot = st.empty()  # diisi belakangan (setelah sel_site dihitung), tapi tampil di atas Divisi

    st.markdown("---")
    st.markdown("### 🏭 Divisi")
    sel_divisi = st.multiselect("Divisi (Mining / Plantation)", list(DIVISI_MAP.keys()), default=list(DIVISI_MAP.keys()))
    sites_allowed_by_divisi = [s for d in sel_divisi for s in DIVISI_MAP.get(d, [])]

    st.markdown("---")
    st.markdown("### 🔎 Filter")

    all_sites_raw = sorted(df_raw["lokasi"].dropna().unique().tolist())
    site_opts = [s for s in all_sites_raw if s in sites_allowed_by_divisi] if sel_divisi else []
    sel_site = st.multiselect("Site / Lokasi", site_opts, default=site_opts)

    # --- Tombol download data maintenance (sudah ada kolom kategori-nya), mengikuti filter Site/Divisi di atas ---
    if not maint_raw.empty:
        maint_dl = maint_raw[maint_raw["lokasi"].isin(sel_site)] if (sel_site and "lokasi" in maint_raw.columns) else maint_raw
        with _download_maint_slot.container():
            st.download_button(
                "⬇️ Download Biaya Maintenance",
                data=maint_dl.to_csv(index=False).encode("utf-8"),
                file_name="data_maintenance.csv", mime="text/csv", use_container_width=True,
            )

    month_opts = [m for m in MONTH_ORDER if m in df_raw["bulan"].unique()]
    sel_month = st.multiselect("Bulan", month_opts, default=month_opts)

    kat_opts = sorted(df_raw["kategori"].dropna().unique().tolist())
    kat_labels = [KATEGORI_LABEL.get(k, k) for k in kat_opts]
    sel_kat_labels = st.multiselect("Kategori Unit", kat_labels, default=kat_labels)
    sel_kat = [k for k in kat_opts if KATEGORI_LABEL.get(k, k) in sel_kat_labels]

    kriteria_scope_df = df_raw[df_raw["lokasi"].isin(sel_site) & df_raw["kategori"].isin(sel_kat)]
    kriteria_opts_raw = sorted(kriteria_scope_df["kriteria_unit"].dropna().unique().tolist()) if "kriteria_unit" in df_raw.columns else []
    has_null_kriteria = ("kriteria_unit" in df_raw.columns) and kriteria_scope_df["kriteria_unit"].isna().any()
    kriteria_opts = kriteria_opts_raw + (["Tidak Diketahui"] if has_null_kriteria else [])
    if kriteria_opts:
        sel_kriteria = st.multiselect("Kriteria Unit (Tarif)", kriteria_opts, default=kriteria_opts)
    else:
        sel_kriteria = []

    unit_opts_df = df_raw[
        df_raw["lokasi"].isin(sel_site) & df_raw["kategori"].isin(sel_kat)
    ][["nama_unit"]].dropna().drop_duplicates().copy()
    unit_opts_df["unit_label"] = unit_opts_df["nama_unit"].apply(_unit_label)
    unit_opts = sorted(unit_opts_df["unit_label"].unique().tolist())
    sel_id_unit = st.multiselect(
        "ID Unit (opsional, kosongkan = semua unit) — ketik ID Unit atau nama unit",
        unit_opts, default=[],
    )

# ---------------------------------------------------------------
# APPLY FILTERS (data utama)
# ---------------------------------------------------------------
df = df_raw[
    df_raw["lokasi"].isin(sel_site) &
    df_raw["bulan"].isin(sel_month) &
    df_raw["kategori"].isin(sel_kat)
].copy()

if "kriteria_unit" in df.columns and kriteria_opts:
    sel_kriteria_actual = [k for k in sel_kriteria if k != "Tidak Diketahui"]
    include_null = "Tidak Diketahui" in sel_kriteria
    mask = df["kriteria_unit"].isin(sel_kriteria_actual)
    if include_null:
        mask = mask | df["kriteria_unit"].isna()
    df = df[mask]

if sel_id_unit:
    df["unit_label"] = df["nama_unit"].apply(_unit_label)
    df = df[df["unit_label"].isin(sel_id_unit)]

maint_df_site_bulan = pd.DataFrame()
if not maint_raw.empty:
    maint_df_site_bulan = maint_raw[
        maint_raw["lokasi"].isin(sel_site) &
        maint_raw["bulan"].isin(sel_month)
    ].copy()
    if sel_id_unit:
        maint_df_site_bulan["unit_label"] = maint_df_site_bulan["nama_unit"].apply(_unit_label)
        maint_df_site_bulan = maint_df_site_bulan[maint_df_site_bulan["unit_label"].isin(sel_id_unit)]

sparepart_df_site_bulan = pd.DataFrame()
if not sparepart_raw.empty:
    sparepart_df_site_bulan = sparepart_raw[
        sparepart_raw["lokasi"].isin(sel_site) &
        sparepart_raw["bulan"].isin(sel_month)
    ].copy()
    if sel_id_unit:
        sparepart_df_site_bulan["unit_label"] = sparepart_df_site_bulan["nama_unit"].apply(_unit_label)
        sparepart_df_site_bulan = sparepart_df_site_bulan[sparepart_df_site_bulan["unit_label"].isin(sel_id_unit)]

sasaran_mutu_df = pd.DataFrame()
if not sasaran_mutu_raw.empty:
    sasaran_mutu_df = sasaran_mutu_raw[
        sasaran_mutu_raw["lokasi"].isin(sel_site) &
        sasaran_mutu_raw["bulan"].isin(sel_month) &
        sasaran_mutu_raw["kategori"].isin(sel_kat) &
        (~sasaran_mutu_raw["unit_sewa"])
    ].copy()
    if sel_id_unit:
        sasaran_mutu_df["unit_label"] = sasaran_mutu_df["nama_unit"].apply(_unit_label)
        sasaran_mutu_df = sasaran_mutu_df[sasaran_mutu_df["unit_label"].isin(sel_id_unit)]

if sel_id_unit:
    st.caption(f"🔗 Dashboard sedang difilter untuk ID Unit: {', '.join(sel_id_unit)}")


def fmt_rp(x):
    if abs(x) >= 1e9:
        return f"Rp {x/1e9:,.2f} M"
    if abs(x) >= 1e6:
        return f"Rp {x/1e6:,.1f} Jt"
    if abs(x) >= 1e3:
        return f"Rp {x/1e3:,.1f} Rb"
    return f"Rp {x:,.0f}"

def achievement(real, budget):
    if budget == 0:
        return None
    return real / budget * 100

# ---------------------------------------------------------------
# HEADER
# ---------------------------------------------------------------
LOGO_PATH = Path(__file__).parent / "logo.png"

def get_logo_base64() -> str:
    if LOGO_PATH.exists():
        with open(LOGO_PATH, "rb") as f:
            return base64.b64encode(f.read()).decode()
    return ""

logo_b64 = get_logo_base64()
logo_html = f'<img class="header-logo" src="data:image/png;base64,{logo_b64}">' if logo_b64 else ""

st.markdown(f"""
<div class="header-banner">
    {logo_html}
    <h1>📊 Dashboard Operational Review</h1>
    <p>PT BUANA KARYA MANDIRI SEJAHTERA (BKMS) &nbsp;•&nbsp; Target vs Realisasi</p>
</div>
""", unsafe_allow_html=True)

if df.empty:
    st.warning("Tidak ada data untuk kombinasi filter yang dipilih. Silakan ubah filter di sidebar.")
    st.stop()

st.caption(f"Site: {', '.join(sel_site) if len(sel_site)<=4 else f'{len(sel_site)} site'} • Bulan: {', '.join(sel_month)}")

# ---------------------------------------------------------------
# HITUNG METRIK UTAMA (dipakai di beberapa bagian + PPTX)
# ---------------------------------------------------------------
tot_pendapatan_r = df["pendapatan_realisasi"].sum()
tot_pendapatan_b = df["pendapatan_budget"].sum()
tot_prestasi_r = df["prestasi_realisasi"].sum()
tot_prestasi_b = df["prestasi_budget"].sum()
tot_biaya_r = df["total_biaya_realisasi"].sum()
tot_biaya_b = df["total_biaya_budget"].sum()
tot_biaya_langsung_r = df["biaya_langsung_realisasi"].sum()
tot_biaya_langsung_b = df["biaya_langsung_budget"].sum()
tot_biaya_tdklangsung_r = df["biaya_tidak_langsung_realisasi"].sum()
tot_biaya_tdklangsung_b = df["biaya_tidak_langsung_budget"].sum()

ach_pendapatan = achievement(tot_pendapatan_r, tot_pendapatan_b)
ach_prestasi = achievement(tot_prestasi_r, tot_prestasi_b)
ach_biaya = achievement(tot_biaya_r, tot_biaya_b)
ach_biaya_langsung = achievement(tot_biaya_langsung_r, tot_biaya_langsung_b)
ach_biaya_tdklangsung = achievement(tot_biaya_tdklangsung_r, tot_biaya_tdklangsung_b)

if "bulan_no" in df.columns and not df.empty:
    _bulan_terakhir_no_global = df["bulan_no"].max()
    df_bulan_terakhir_global = df[df["bulan_no"] == _bulan_terakhir_no_global]
else:
    df_bulan_terakhir_global = df

target_populasi = df_bulan_terakhir_global.loc[df_bulan_terakhir_global["pendapatan_budget"] > 0, "nama_unit"].nunique()
realisasi_populasi = df_bulan_terakhir_global.loc[df_bulan_terakhir_global["pendapatan_realisasi"] > 0, "nama_unit"].nunique()
pct_populasi = (realisasi_populasi / target_populasi * 100) if target_populasi else None

# ---------------------------------------------------------------
# POWERPOINT EXPORT — didesain mengikuti gaya "Tinjauan Manajemen" (RTM):
# background terang, header bar navy, kartu KPI ikon+pill status, tabel
# dengan indikator warna, dan kotak analisis.
# ---------------------------------------------------------------
def build_pptx(data, maint_data, sparepart_data, site_list, month_list, kat_list, sasaran_mutu_data=None) -> bytes:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION, XL_MARKER_STYLE
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.dml import MSO_LINE_DASH_STYLE
    import io as _io

    # Khusus site TANJUNG: gabungkan kategori TR ke AB (tidak dipisah AB/TR di seluruh chart PPTX)
    data = data.copy()
    data.loc[data["lokasi"] == "TANJUNG", "kategori"] = "AB"
    if sasaran_mutu_data is not None and not sasaran_mutu_data.empty:
        sasaran_mutu_data = sasaran_mutu_data.copy()
        sasaran_mutu_data.loc[sasaran_mutu_data["lokasi"] == "TANJUNG", "kategori"] = "AB"
        sasaran_mutu_data.loc[sasaran_mutu_data["lokasi"] == "TANJUNG", "Jenis_Sarmut"] = "Sarmut Kelompok Alat Berat"

    # Singkatan nama site dipakai konsisten di semua chart yg padat kategori
    SITE_ABBR = {"SUNGAI DANAU": "S.DANAU", "BUHUT LHL": "B.LHL", "TANJUNG": "TANJUNG",
                 "BUHUT": "BUHUT", "KUMAI": "KUMAI", "AMPAH": "AMPAH"}

    # Korelasi Jenis Unit -> Jenis Sarmut (mis. Arm Roll -> Dump Truck)
    JENIS_UNIT_TO_SARMUT = {
        "DT": "Dump Truck", "DT FM 260 JD": "Dump Truck", "DT Hino": "Dump Truck",
        "DT Howo": "Dump Truck", "Arm Roll": "Dump Truck", "Trailler": "Dump Truck",
        "Isuzu Elf": "Dump Truck",
        "Crawler": "Crawler",
        "Excavator": "Excavator", "Excavator All in": "Excavator", "Excavator Braaker 20 Ton": "Excavator",
        "Excavator Breaker 30 Ton": "Excavator", "Excavator Bucket": "Excavator",
        "Excavator Long arm": "Excavator", "Excavator Mini": "Excavator", "Excavator ZX-130": "Excavator",
        "FT": "FT",
        "Bulldozer": "Unit Support", "Compactor": "Unit Support", "Grader": "Unit Support",
        "TLB": "Unit Support", "Backhoe Loader": "Unit Support", "Wheel Loader": "Unit Support",
        "Pompa Air": "Unit Support", "Fuel Tank": "Unit Support",
        "Truk Tangki": "Tangki CPO", "Tangki Air": "Tangki CPO",
    }

    def map_jenis_sarmut(jenis_unit):
        return JENIS_UNIT_TO_SARMUT.get(jenis_unit, "Alat Berat")

    # Deteksi divisi (Mining/Plantation) berdasarkan site yang difilter
    MINING_SITES_PPTX = {"TANJUNG", "BUHUT", "BUHUT LHL"}
    PLANTATION_SITES_PPTX = {"SUNGAI DANAU", "KUMAI"}
    site_set = set(site_list) if site_list else set()
    if site_set and site_set.issubset(MINING_SITES_PPTX):
        divisi_label = " · MINING"
    elif site_set and site_set.issubset(PLANTATION_SITES_PPTX):
        divisi_label = " · PLANTATION"
    else:
        divisi_label = ""

    # --- Palet warna (mengikuti gaya laporan RTM: terang, header navy) ---
    NAVY = RGBColor(0x1B, 0x25, 0x4B)
    NAVY_DARK = RGBColor(0x12, 0x18, 0x35)
    BG_LIGHT = RGBColor(0xF2, 0xF4, 0xF9)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    TEAL = RGBColor(0x17, 0xA9, 0xC7)
    GREEN = RGBColor(0x21, 0xA1, 0x66)
    GREEN_BG = RGBColor(0xDF, 0xF5, 0xE9)
    RED = RGBColor(0xE0, 0x4A, 0x3D)
    RED_BG = RGBColor(0xFC, 0xE4, 0xE1)
    GOLD = RGBColor(0xC9, 0x8A, 0x1E)
    GOLD_BG = RGBColor(0xFB, 0xEE, 0xD8)
    TEXT_DARK = RGBColor(0x21, 0x29, 0x37)
    TEXT_MUTED = RGBColor(0x71, 0x78, 0x86)
    BORDER = RGBColor(0xE1, 0xE4, 0xEC)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    def add_bg(slide, color):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid(); bg.fill.fore_color.rgb = color
        bg.line.fill.background(); bg.shadow.inherit = False
        slide.shapes._spTree.remove(bg._element)
        slide.shapes._spTree.insert(2, bg._element)
        return bg

    def add_content_slide(title, subtitle_right=""):
        s = prs.slides.add_slide(blank)
        add_bg(s, BG_LIGHT)
        header = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.95))
        header.fill.solid(); header.fill.fore_color.rgb = NAVY
        header.line.fill.background(); header.shadow.inherit = False
        tb = s.shapes.add_textbox(Inches(0.5), Inches(0.18), Inches(9), Inches(0.6))
        p = tb.text_frame.paragraphs[0]
        r = p.add_run(); r.text = title
        r.font.size = Pt(20); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = "Calibri"
        if subtitle_right:
            tb2 = s.shapes.add_textbox(Inches(9.5), Inches(0.28), Inches(3.4), Inches(0.4))
            p2 = tb2.text_frame.paragraphs[0]
            p2.alignment = PP_ALIGN.RIGHT
            r2 = p2.add_run(); r2.text = subtitle_right
            r2.font.size = Pt(11); r2.font.color.rgb = RGBColor(0xC9, 0xCF, 0xE0); r2.font.name = "Calibri"
        return s

    def add_textbox(slide, left, top, width, height, text, size=14, bold=False,
                     color=TEXT_DARK, align=PP_ALIGN.LEFT, italic=False):
        tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        r = p.add_run(); r.text = text
        r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
        r.font.color.rgb = color; r.font.name = "Calibri"
        return tb

    def pill_colors(is_good):
        return (GREEN_BG, GREEN) if is_good else (RED_BG, RED)

    def add_soft_shadow(shape, blur=90000, dist=22000, alpha=22000):
        """Tambahkan efek bayangan lembut (drop shadow) custom via manipulasi XML,
        karena python-pptx tidak menyediakan API tingkat tinggi untuk ini."""
        from pptx.oxml.ns import qn as _qn_shadow
        spPr = shape._element.spPr
        existing = spPr.find(_qn_shadow('a:effectLst'))
        if existing is not None:
            spPr.remove(existing)
        effectLst = spPr.makeelement(_qn_shadow('a:effectLst'), {})
        outerShdw = spPr.makeelement(_qn_shadow('a:outerShdw'), {
            'blurRad': str(blur), 'dist': str(dist), 'dir': '5400000', 'rotWithShape': '0'
        })
        srgbClr = spPr.makeelement(_qn_shadow('a:srgbClr'), {'val': '1A2744'})
        alpha_el = spPr.makeelement(_qn_shadow('a:alpha'), {'val': str(alpha)})
        srgbClr.append(alpha_el)
        outerShdw.append(srgbClr)
        effectLst.append(outerShdw)
        spPr.append(effectLst)

    def add_kpi_card(slide, left, top, width, height, icon_txt, icon_color, accent_color,
                      label, value, sub_text, pill_text, pill_good):
        # accent strip (top border)
        strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.07))
        strip.fill.solid(); strip.fill.fore_color.rgb = accent_color
        strip.line.fill.background(); strip.shadow.inherit = False
        # card body
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top + 0.07), Inches(width), Inches(height - 0.07))
        card.adjustments[0] = 0.045
        card.fill.solid(); card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = BORDER; card.line.width = Pt(0.75)
        card.shadow.inherit = False
        add_soft_shadow(card)
        # icon circle (ukuran diperbesar sedikit agar lebih menonjol; sedikit lebih kecil kalau kartu sempit)
        narrow_pre = width < 2.6
        icon_size = 0.46 if narrow_pre else 0.56
        icon_left = left + 0.2 if narrow_pre else left + 0.25
        circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(icon_left), Inches(top + 0.25), Inches(icon_size), Inches(icon_size))
        circ.fill.solid(); circ.fill.fore_color.rgb = icon_color
        circ.line.fill.background(); circ.shadow.inherit = False
        ic_tf = circ.text_frame; ic_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        ic_tf.margin_left = 0; ic_tf.margin_right = 0
        icp = ic_tf.paragraphs[0]; icp.alignment = PP_ALIGN.CENTER
        icr = icp.add_run(); icr.text = icon_txt
        icr.font.size = Pt(14 if narrow_pre else 17); icr.font.bold = True; icr.font.color.rgb = WHITE; icr.font.name = EMOJI_FONT
        # Skala ukuran font & posisi menyesuaikan lebar kartu (supaya tetap muat kalau kartu dibuat sempit, mis. 5 kartu sejajar)
        narrow = narrow_pre
        label_size = 9.5 if narrow else 11.5
        value_size = 19 if narrow else 23
        sub_size = 8.5 if narrow else 10
        pill_size = 9 if narrow else 10.5
        label_h = 0.55 if narrow else 0.4
        # label
        add_textbox(slide, left + 1.0, top + 0.24, width - 1.15, label_h, label, size=label_size, bold=True, color=TEXT_MUTED)
        # value (posisi proporsional thd tinggi kartu, agar tidak tumpang tindih di kartu pendek)
        value_top = top + (0.78 if narrow else 0.66)
        add_textbox(slide, left + 0.25, value_top, width - 0.5, 0.5, value, size=value_size, bold=True, color=TEXT_DARK)
        # sub text (target/budget)
        if sub_text:
            add_textbox(slide, left + 0.25, value_top + (0.38 if narrow else 0.42), width - 0.5, 0.3, sub_text, size=sub_size, color=TEXT_MUTED)
        # pill (selalu menempel ke bawah kartu)
        pbg, ptxt = pill_colors(pill_good)
        pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left + 0.25), Inches(top + height - 0.5), Inches(width - 0.5), Inches(0.35))
        pill.adjustments[0] = 0.5
        pill.fill.solid(); pill.fill.fore_color.rgb = pbg
        pill.line.fill.background(); pill.shadow.inherit = False
        ptf = pill.text_frame; ptf.vertical_anchor = MSO_ANCHOR.MIDDLE
        pp = ptf.paragraphs[0]; pp.alignment = PP_ALIGN.CENTER
        pr = pp.add_run(); pr.text = pill_text
        pr.font.size = Pt(pill_size); pr.font.bold = True; pr.font.color.rgb = ptxt

    def add_card_panel(slide, left, top, width, height, accent_color=None):
        card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        card.fill.solid(); card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = BORDER; card.line.width = Pt(0.75)
        card.shadow.inherit = False
        add_soft_shadow(card, blur=70000, dist=18000, alpha=18000)
        if accent_color:
            strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.06))
            strip.fill.solid(); strip.fill.fore_color.rgb = accent_color
            strip.line.fill.background(); strip.shadow.inherit = False
        return card

    def add_note_callout(slide, left, top, width, height, icon, text, text_color=RED, size=11):
        tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = tb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]
        r_icon = p.add_run(); r_icon.text = f"{icon} "
        r_icon.font.size = Pt(size); r_icon.font.bold = True; r_icon.font.color.rgb = text_color; r_icon.font.name = EMOJI_FONT
        r = p.add_run(); r.text = text
        r.font.size = Pt(size); r.font.bold = True; r.font.italic = True
        r.font.color.rgb = text_color; r.font.name = "Calibri"
        return tb

    def add_finding_box(slide, left, top, width, height, icon, text, bg_color, border_color, text_color):
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        box.adjustments[0] = min(0.12, 0.35 / height)
        box.fill.solid(); box.fill.fore_color.rgb = bg_color
        box.line.color.rgb = border_color; box.line.width = Pt(1.25)
        box.shadow.inherit = False
        # Ikon dalam lingkaran (badge) di kiri, konsisten dgn elemen lain, supaya lebih menonjol
        icon_d = min(0.42, height - 0.16)
        icon_y = top + height / 2 - icon_d / 2
        circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(left + 0.14), Inches(icon_y), Inches(icon_d), Inches(icon_d))
        circ.fill.solid(); circ.fill.fore_color.rgb = border_color
        circ.line.fill.background(); circ.shadow.inherit = False
        ictf = circ.text_frame; ictf.vertical_anchor = MSO_ANCHOR.MIDDLE
        ictf.margin_left = 0; ictf.margin_right = 0
        icp = ictf.paragraphs[0]; icp.alignment = PP_ALIGN.CENTER
        icr = icp.add_run(); icr.text = icon
        icr.font.size = Pt(max(10, icon_d * 22)); icr.font.color.rgb = WHITE; icr.font.name = EMOJI_FONT
        text_left = left + 0.14 + icon_d + 0.14
        tf = slide.shapes.add_textbox(Inches(text_left), Inches(top), Inches(width - (text_left - left) - 0.15), Inches(height)).text_frame
        tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        r = p.add_run(); r.text = text
        r.font.size = Pt(10.5); r.font.bold = True; r.font.color.rgb = text_color; r.font.name = "Calibri"
        return box


    def add_status_banner(slide, left, top, width, height, icon, text, bg_color, border_color, text_color):
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        box.adjustments[0] = 0.15
        box.fill.solid(); box.fill.fore_color.rgb = bg_color
        box.line.color.rgb = border_color; box.line.width = Pt(1)
        box.shadow.inherit = False
        strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(0.07), Inches(height))
        strip.fill.solid(); strip.fill.fore_color.rgb = border_color
        strip.line.fill.background(); strip.shadow.inherit = False
        tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.2)
        p = tf.paragraphs[0]
        r_icon = p.add_run(); r_icon.text = f"{icon}  "
        r_icon.font.size = Pt(12.5); r_icon.font.bold = True; r_icon.font.color.rgb = text_color; r_icon.font.name = EMOJI_FONT
        r = p.add_run(); r.text = text
        r.font.size = Pt(12.5); r.font.bold = True; r.font.color.rgb = text_color; r.font.name = "Calibri"
        return box

    def add_table(slide, left, top, width, height, headers, rows, status_col=None, col_widths=None,
                  fill_badge=False, font_size=11.5, header_size=12, status_row=None):
        n_rows = len(rows) + 1
        n_cols = len(headers)
        gframe = slide.shapes.add_table(n_rows, n_cols, Inches(left), Inches(top), Inches(width), Inches(height))
        table = gframe.table
        if col_widths:
            total = sum(col_widths)
            for i, w in enumerate(col_widths):
                table.columns[i].width = Inches(width * w / total)
        status_cols = [status_col] if isinstance(status_col, int) else (status_col or [])
        status_rows = [status_row] if isinstance(status_row, int) else (status_row or [])
        for j, h in enumerate(headers):
            cell = table.cell(0, j)
            cell.text = h
            cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
            p = cell.text_frame.paragraphs[0]
            p.runs[0].font.bold = True; p.runs[0].font.size = Pt(header_size); p.runs[0].font.color.rgb = WHITE
            p.runs[0].font.name = "Calibri"
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        for i, row in enumerate(rows, start=1):
            for j, val in enumerate(row):
                cell = table.cell(i, j)
                text = str(val)
                is_status = (j in status_cols) or (i in status_rows and j > 0)
                alt_bg = WHITE if i % 2 == 1 else RGBColor(0xF5, 0xF7, 0xFB)
                is_na = text.startswith("N/A") or text == "-"
                is_good = text.startswith("✓")
                display_text = text.lstrip("✓✗").strip()
                if is_status and fill_badge:
                    if is_na:
                        cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0x9C, 0xA3, 0xAF)
                        txt_color = WHITE
                    elif is_good:
                        cell.fill.solid(); cell.fill.fore_color.rgb = alt_bg
                        txt_color = GREEN
                    else:
                        cell.fill.solid(); cell.fill.fore_color.rgb = RED
                        txt_color = WHITE
                elif is_status:
                    cell.fill.solid(); cell.fill.fore_color.rgb = alt_bg
                    txt_color = RGBColor(0x9C, 0xA3, 0xAF) if is_na else (GREEN if is_good else RED)
                else:
                    cell.fill.solid(); cell.fill.fore_color.rgb = alt_bg
                    txt_color = TEXT_DARK
                cell.text_frame.paragraphs[0].text = ""
                p = cell.text_frame.paragraphs[0]
                r = p.add_run(); r.text = display_text
                r.font.size = Pt(font_size)
                r.font.color.rgb = txt_color
                r.font.bold = is_status
                r.font.name = "Calibri"
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                cell.margin_left = Inches(0.08); cell.margin_right = Inches(0.08)
        return table

    def add_insight_box(slide, left, top, width, height, title, bullets, border_color=TEAL, title_color=None):
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        box.adjustments[0] = 0.03
        box.fill.solid(); box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = border_color; box.line.width = Pt(1.25)
        box.shadow.inherit = False
        left_strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(0.07), Inches(height))
        left_strip.fill.solid(); left_strip.fill.fore_color.rgb = border_color
        left_strip.line.fill.background(); left_strip.shadow.inherit = False
        tb = slide.shapes.add_textbox(Inches(left + 0.3), Inches(top + 0.15), Inches(width - 0.55), Inches(height - 0.3))
        tf = tb.text_frame; tf.word_wrap = True
        p0 = tf.paragraphs[0]
        r0 = p0.add_run(); r0.text = title
        r0.font.size = Pt(13); r0.font.bold = True; r0.font.color.rgb = (title_color or border_color)
        p0.space_after = Pt(8)
        for b in bullets:
            p = tf.add_paragraph()
            r = p.add_run(); r.text = f"•  {b}"
            r.font.size = Pt(11); r.font.color.rgb = TEXT_DARK
            p.space_after = Pt(6)

    def style_chart_light(chart, legend=True, legend_pos=None):
        chart.has_title = False
        chart.has_legend = legend
        if legend:
            chart.legend.position = legend_pos or XL_LEGEND_POSITION.BOTTOM
            chart.legend.include_in_layout = False
            chart.legend.font.color.rgb = TEXT_DARK
            chart.legend.font.size = Pt(10.5)
            chart.legend.font.bold = True
            chart.legend.font.name = "Calibri"
        cat_ax = chart.category_axis
        cat_ax.tick_labels.font.color.rgb = TEXT_DARK
        cat_ax.tick_labels.font.size = Pt(10)
        cat_ax.tick_labels.font.bold = True
        cat_ax.tick_labels.font.name = "Calibri"
        cat_ax.format.line.color.rgb = BORDER
        val_ax = chart.value_axis
        val_ax.tick_labels.font.color.rgb = TEXT_DARK
        val_ax.tick_labels.font.size = Pt(9.5)
        val_ax.tick_labels.font.name = "Calibri"
        val_ax.tick_labels.number_format = '#,##0'
        val_ax.tick_labels.number_format_is_linked = False
        val_ax.format.line.color.rgb = BORDER
        val_ax.has_major_gridlines = True
        val_ax.major_gridlines.format.line.color.rgb = BORDER

    def add_panel_header(slide, left, top, width, text, height=0.42):
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        bar.fill.solid(); bar.fill.fore_color.rgb = NAVY
        bar.line.fill.background(); bar.shadow.inherit = False
        # Aksen tipis di bawah header supaya ada pemisah visual halus dgn konten panel
        accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top + height - 0.035), Inches(width), Inches(0.035))
        accent.fill.solid(); accent.fill.fore_color.rgb = GOLD
        accent.line.fill.background(); accent.shadow.inherit = False
        tf = bar.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.15); tf.margin_right = Inches(0.1)
        p = tf.paragraphs[0]
        # Pisahkan ikon (kata pertama) dari teks, ikon dirender sedikit lebih besar supaya lebih menonjol
        parts = text.split(" ", 1)
        if len(parts) == 2 and len(parts[0]) <= 2:
            r_icon = p.add_run(); r_icon.text = parts[0] + "  "
            r_icon.font.size = Pt(15); r_icon.font.bold = True; r_icon.font.color.rgb = WHITE; r_icon.font.name = EMOJI_FONT
            r_txt = p.add_run(); r_txt.text = parts[1]
            r_txt.font.size = Pt(12.5); r_txt.font.bold = True; r_txt.font.color.rgb = WHITE; r_txt.font.name = "Calibri"
        else:
            r = p.add_run(); r.text = text
            r.font.size = Pt(12.5); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = "Calibri"
        return bar

    def ach_txt_pct(real, budget):
        if budget == 0:
            return None
        return real / budget * 100

    import datetime as _dt

    period = sorted(month_list, key=lambda m: MONTH_ORDER.index(m) if m in MONTH_ORDER else 99)[-1] if month_list else "-"
    site_txt = ", ".join(site_list) if len(site_list) <= 6 else f"{len(site_list)} site"
    kat_txt = ", ".join([KATEGORI_LABEL.get(k, k) for k in kat_list])
    _bulan_id = ["Januari","Februari","Maret","April","Mei","Juni","Juli","Agustus","September","Oktober","November","Desember"]
    _now = _dt.datetime.now()
    tgl_laporan = f"{_bulan_id[_now.month-1]} {_now.year}"

    def render_6_slides(data, sasaran_mutu_data, snum1, snum2, snum3, snum4, snum5, snum6, kat_suffix):
        r_ = data["pendapatan_realisasi"].sum(); b_ = data["pendapatan_budget"].sum()
        pr_ = data["prestasi_realisasi"].sum(); pb_ = data["prestasi_budget"].sum()
        bl_r_raw = data["biaya_langsung_realisasi"].sum(); bl_b_raw = data["biaya_langsung_budget"].sum()
        btl_r_raw = data["biaya_tidak_langsung_realisasi"].sum(); btl_b_raw = data["biaya_tidak_langsung_budget"].sum()
        bl_r = (bl_r_raw / pr_) if pr_ else None
        bl_b = (bl_b_raw / pb_) if pb_ else None
        btl_r = (btl_r_raw / pr_) if pr_ else None
        btl_b = (btl_b_raw / pb_) if pb_ else None
        ach_r = ach_txt_pct(r_, b_); ach_p = ach_txt_pct(pr_, pb_)
        ach_bl = ach_txt_pct(bl_r, bl_b) if (bl_r is not None and bl_b) else None
        ach_btl = ach_txt_pct(btl_r, btl_b) if (btl_r is not None and btl_b) else None

        if sasaran_mutu_data is None:
            sasaran_mutu_data = pd.DataFrame()

        def klasifikasi_satuan_lokal(row):
            lok, kat = row["lokasi"], row["kategori"]
            if lok == "BUHUT LHL":
                return "Ton"
            if lok in ("SUNGAI DANAU", "KUMAI"):
                if kat == "AB":
                    return "HM"
                if kat == "TR":
                    return "KM"
                return "HM"
            return "HM"

        data["satuan_lokal"] = data.apply(klasifikasi_satuan_lokal, axis=1)

        avg_avail_r = sasaran_mutu_data["availability_pct"].mean() if not sasaran_mutu_data.empty else None
        avg_avail_t = sasaran_mutu_data["availability_target"].mean() if not sasaran_mutu_data.empty else None
        avg_util_r = sasaran_mutu_data["utilisasi_pct"].mean() if not sasaran_mutu_data.empty else None
        avg_util_t = sasaran_mutu_data["utilisasi_target"].mean() if not sasaran_mutu_data.empty else None
        ach_avail = ach_txt_pct(avg_avail_r, avg_avail_t) if (avg_avail_r is not None and avg_avail_t) else None
        ach_util = ach_txt_pct(avg_util_r, avg_util_t) if (avg_util_r is not None and avg_util_t) else None
        prestasi_r_kpi = data["prestasi_realisasi"].sum() if "prestasi_realisasi" in data.columns else None
        prestasi_b_kpi = data["prestasi_budget"].sum() if "prestasi_budget" in data.columns else None
        ach_prestasi_kpi = ach_txt_pct(prestasi_r_kpi, prestasi_b_kpi) if (prestasi_r_kpi is not None and prestasi_b_kpi) else None

        # ================= SLIDE 1: KPI DASHBOARD — PERFORMANCE KESELURUHAN =================
        s = add_content_slide(f"KPI DASHBOARD — Performance Keseluruhan s/d {period}", f"Ringkasan Kinerja · {snum1}{divisi_label}{kat_suffix}")

        # --- Siapkan data chart: % Capaian Utilisasi & % Capaian Prestasi per Site & Jenis Unit ---
        # Dipisah jadi 2 CHART TERPISAH: Floating Tarif & Tarif Tetap (berdasarkan kolom kriteria_unit)
        au_rows = []

        kriteria_unit_lookup_grp = data.dropna(subset=["kriteria_unit"]).groupby(
            ["lokasi", "kategori", "jenis_unit"])["kriteria_unit"].agg(lambda s: s.mode().iloc[0] if not s.mode().empty else None).to_dict()

        data_k = data.copy()

        if not sasaran_mutu_data.empty:
            sm_kpi = sasaran_mutu_data.copy()
            sm_kpi = sm_kpi.dropna(subset=["jenis_unit"])

            au_tbl = sm_kpi.groupby(["lokasi", "kategori", "jenis_unit"], as_index=False).agg(
                avail_r=("availability_pct", "mean"), avail_t=("availability_target", "mean"),
                util_r=("utilisasi_pct", "mean"), util_t=("utilisasi_target", "mean"),
            )
            au_tbl["site_short"] = au_tbl["lokasi"].map(SITE_ABBR).fillna(au_tbl["lokasi"])
            au_tbl = au_tbl.sort_values(["lokasi", "kategori", "jenis_unit"])

            prestasi_unit = data_k.groupby(["lokasi", "kategori", "jenis_unit"], as_index=False).agg(
                prestasi_r=("prestasi_realisasi", "sum"), prestasi_b=("prestasi_budget", "sum"))
            prestasi_unit["cap"] = prestasi_unit.apply(lambda r: (r["prestasi_r"] / r["prestasi_b"] * 100) if r["prestasi_b"] else None, axis=1)
            prestasi_lookup_unit = {(r["lokasi"], r["kategori"], r["jenis_unit"]): r["cap"] for _, r in prestasi_unit.iterrows()}

            # Fallback ke level site+kategori kalau kombinasi jenis_unit spesifik tidak ada datanya
            prestasi_sk1 = data.groupby(["lokasi", "kategori"], as_index=False).agg(
                prestasi_r=("prestasi_realisasi", "sum"), prestasi_b=("prestasi_budget", "sum"))
            prestasi_sk1["cap"] = prestasi_sk1.apply(lambda r: (r["prestasi_r"] / r["prestasi_b"] * 100) if r["prestasi_b"] else None, axis=1)
            prestasi_lookup = {(r["lokasi"], r["kategori"]): r["cap"] for _, r in prestasi_sk1.iterrows()}

            for _, r in au_tbl.iterrows():
                util_cap = (r["util_r"] / r["util_t"] * 100) if r["util_t"] else None
                prestasi_cap = prestasi_lookup_unit.get((r["lokasi"], r["kategori"], r["jenis_unit"]))
                if prestasi_cap is None:
                    prestasi_cap = prestasi_lookup.get((r["lokasi"], r["kategori"]))
                kriteria = kriteria_unit_lookup_grp.get((r["lokasi"], r["kategori"], r["jenis_unit"]))
                au_rows.append({"label": f"{r['site_short']} — {r['jenis_unit']}", "util_cap": util_cap,
                                 "prestasi_cap": prestasi_cap, "kriteria_unit": kriteria})

        au_rows_floating = [r for r in au_rows if r["kriteria_unit"] == "Floating Tarif"]
        au_rows_tetap = [r for r in au_rows if r["kriteria_unit"] == "Tarif Tetap"]
        au_rows_other = [r for r in au_rows if r["kriteria_unit"] not in ("Floating Tarif", "Tarif Tetap")]
        au_rows_floating = au_rows_floating + au_rows_other  # unit tanpa info kriteria digabung ke Floating (default)

        # --- Kartu ringkasan mini (ringkasan cepat keseluruhan, lengkap dgn Budget & Capaian) ---
        mini_w, mini_h, mini_gap, mini_y = 2.32, 1.95, 0.19, 1.05
        add_kpi_card(s, 0.4, mini_y, mini_w, mini_h, "📈", RGBColor(0x2E, 0x6D, 0xB4), GREEN if (ach_prestasi_kpi is not None and ach_prestasi_kpi >= 100) else RED,
                     "Capaian Prestasi", (f"{ach_prestasi_kpi:.1f}%" if ach_prestasi_kpi is not None else "-"),
                     "Target: 100.0%" if prestasi_r_kpi is not None else "Data tidak tersedia",
                     (f"✓ {ach_prestasi_kpi:.1f}% — Tercapai" if ach_prestasi_kpi is not None and ach_prestasi_kpi >= 100 else (f"✗ {ach_prestasi_kpi:.1f}% — Belum Tercapai" if ach_prestasi_kpi is not None else "Data tidak tersedia")),
                     ach_prestasi_kpi is not None and ach_prestasi_kpi >= 100)
        add_kpi_card(s, 0.4 + (mini_w + mini_gap), mini_y, mini_w, mini_h, "🎯", GOLD, GOLD if (ach_util is not None and ach_util < 100) else GREEN,
                     "Avg Utilisasi", (f"{avg_util_r:.1f}%" if avg_util_r is not None else "-"),
                     f"Target: {avg_util_t:.1f}%" if avg_util_t is not None else "Target: -",
                     (f"✓ {ach_util:.1f}% dari Target" if ach_util is not None and ach_util >= 100 else (f"✗ {ach_util:.1f}% dari Target" if ach_util is not None else "Data tidak tersedia")),
                     ach_util is not None and ach_util >= 100)
        add_kpi_card(s, 0.4 + 2 * (mini_w + mini_gap), mini_y, mini_w, mini_h, "⚙", TEAL, TEAL if (ach_avail is not None and ach_avail < 100) else GREEN,
                     "Avg Availability", (f"{avg_avail_r:.1f}%" if avg_avail_r is not None else "-"),
                     f"Target: {avg_avail_t:.1f}%" if avg_avail_t is not None else "Target: -",
                     (f"✓ {ach_avail:.1f}% dari Target" if ach_avail is not None and ach_avail >= 100 else (f"✗ {ach_avail:.1f}% dari Target" if ach_avail is not None else "Data tidak tersedia")),
                     ach_avail is not None and ach_avail >= 100)
        add_kpi_card(s, 0.4 + 3 * (mini_w + mini_gap), mini_y, mini_w, mini_h, "💰", GREEN, GREEN if (ach_bl is not None and ach_bl <= 100) else RED,
                     "Biaya Langsung / Prestasi", (fmt_rp(bl_r) if bl_r is not None else "-"),
                     f"Budget: {fmt_rp(bl_b)}" if bl_b is not None else "Budget: -",
                     (f"✓ {ach_bl:.1f}% — Under Budget" if ach_bl is not None and ach_bl <= 100 else (f"✗ {ach_bl:.1f}% — Over Budget" if ach_bl is not None else "Target = 0")),
                     ach_bl is not None and ach_bl <= 100)
        add_kpi_card(s, 0.4 + 4 * (mini_w + mini_gap), mini_y, mini_w, mini_h, "🧾", RGBColor(0x8E, 0x6B, 0xC9), GREEN if (ach_btl is not None and ach_btl <= 100) else RED,
                     "Biaya T.Langsung / Prestasi", (fmt_rp(btl_r) if btl_r is not None else "-"),
                     f"Budget: {fmt_rp(btl_b)}" if btl_b is not None else "Budget: -",
                     (f"✓ {ach_btl:.1f}% — Under Budget" if ach_btl is not None and ach_btl <= 100 else (f"✗ {ach_btl:.1f}% — Over Budget" if ach_btl is not None else "Target = 0")),
                     ach_btl is not None and ach_btl <= 100)

        def _draw_util_chart(slide, rows, top, height, title):
            add_card_panel(slide, 0.45, top, 12.35, height)
            add_panel_header(slide, 0.45, top, 12.35, title, height=0.34)
            chart_top_x = top + 0.4
            chart_h_x = height - 0.45
            if rows:
                n_x = len(rows)
                # Label satuan Utilisasi menyesuaikan kategori: Transportasi -> Hari, Alat Berat -> HM.
                # Kalau datanya campuran (AB & TR sekaligus, tidak dipisah), label dikosongkan (netral).
                kat_set_util = set(data["kategori"].dropna().unique())
                if kat_set_util == {"TR"}:
                    util_label_x = "% Capaian Utilisasi (Hari)"
                elif kat_set_util == {"AB"}:
                    util_label_x = "% Capaian Utilisasi (HM)"
                else:
                    util_label_x = "% Capaian Utilisasi"
                cd_x = CategoryChartData()
                cd_x.categories = [r["label"] for r in rows]
                cd_x.add_series("% Capaian Prestasi", tuple(round(r["prestasi_cap"], 1) if r["prestasi_cap"] is not None else 0 for r in rows))
                cd_x.add_series(util_label_x, tuple(round(r["util_cap"], 1) if r["util_cap"] is not None else 0 for r in rows))
                gframe_x = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.6), Inches(chart_top_x), Inches(12.2), Inches(chart_h_x), cd_x)
                chart_x = gframe_x.chart
                PRESTASI_COLOR = RGBColor(0x2E, 0x6D, 0xB4)  # disamakan dgn warna ikon kartu KPI "Capaian Prestasi"
                chart_x.series[0].format.fill.solid(); chart_x.series[0].format.fill.fore_color.rgb = PRESTASI_COLOR
                chart_x.series[1].format.fill.solid(); chart_x.series[1].format.fill.fore_color.rgb = GOLD
                chart_x.has_title = False
                plot_x = chart_x.plots[0]
                plot_x.gap_width = 60
                plot_x.has_data_labels = True
                dls_x = plot_x.data_labels
                dls_x.number_format = '0"%"'; dls_x.number_format_is_linked = False
                dls_x.font.size = Pt(6.5); dls_x.font.bold = True; dls_x.font.color.rgb = TEXT_DARK; dls_x.font.name = "Calibri"
                dls_x.position = XL_LABEL_POSITION.OUTSIDE_END
                # Label angka Prestasi ditampilkan utk SEMUA nilai (termasuk yg >=100%), diwarnai merah kalau <100%
                for i_pt_x, pt_x in enumerate(chart_x.series[0].points):
                    prestasi_val_x = rows[i_pt_x]["prestasi_cap"]
                    dl_x = pt_x.data_label
                    dl_x.has_text_frame = True
                    if prestasi_val_x is not None:
                        dl_x.text_frame.text = f"{prestasi_val_x:.0f}%"
                        r0_x = dl_x.text_frame.paragraphs[0].runs[0]
                        r0_x.font.size = Pt(6.5); r0_x.font.bold = True; r0_x.font.name = "Calibri"
                        r0_x.font.color.rgb = RED if prestasi_val_x < 100 else TEXT_DARK
                    else:
                        dl_x.text_frame.text = "-"
                        r0_x = dl_x.text_frame.paragraphs[0].runs[0]
                        r0_x.font.size = Pt(6.5); r0_x.font.bold = True; r0_x.font.name = "Calibri"; r0_x.font.color.rgb = TEXT_MUTED
                style_chart_light(chart_x, legend=True, legend_pos=XL_LEGEND_POSITION.BOTTOM)
                cat_font_x = 8.5 if n_x <= 8 else (7 if n_x <= 14 else (6 if n_x <= 22 else 5.3))
                chart_x.category_axis.tick_labels.font.size = Pt(cat_font_x)
                chart_x.value_axis.tick_labels.font.size = Pt(cat_font_x)
                chart_x.value_axis.tick_labels.number_format = '0"%"'
                chart_x.value_axis.tick_labels.number_format_is_linked = False
            else:
                add_textbox(slide, 0.6, chart_top_x + 0.1, 12.0, 0.4, "Tidak ada data untuk kategori ini.", size=10, italic=True, color=TEXT_MUTED)

        # --- Chart Floating Tarif saja (Tarif Tetap dihilangkan krn tidak ada tracking Utilisasi/Availability yg berarti) ---
        panel_top = mini_y + mini_h + 0.15
        panel_bottom = 7.3
        total_h = panel_bottom - panel_top
        h_floating = total_h * 0.8
        h_finding = total_h - h_floating - 0.15

        _draw_util_chart(s, au_rows_floating, panel_top, h_floating, "🟢 Capaian Utilisasi & Prestasi — per Site & Jenis Unit")

        # --- Analisa: unit dgn gap pendapatan (realisasi - budget) paling minus, dikaitkan dgn capaian utilisasinya ---
        note_top_au = panel_top + h_floating + 0.1
        note_h_au = panel_bottom - note_top_au
        gap_unit = data_k.groupby(["lokasi", "kategori", "jenis_unit"], as_index=False).agg(
            pend_r=("pendapatan_realisasi", "sum"), pend_b=("pendapatan_budget", "sum"))
        gap_unit = gap_unit[(gap_unit["pend_r"] > 0) | (gap_unit["pend_b"] > 0)].copy()
        gap_unit["gap"] = gap_unit["pend_r"] - gap_unit["pend_b"]
        gap_unit["site_short"] = gap_unit["lokasi"].map(SITE_ABBR).fillna(gap_unit["lokasi"])
        gap_unit["label"] = gap_unit["site_short"] + " — " + gap_unit["jenis_unit"]
        util_cap_lookup = {r["label"]: r["prestasi_cap"] for r in au_rows}
        gap_unit["prestasi_cap"] = gap_unit["label"].map(util_cap_lookup)
        gap_unit_neg = gap_unit[gap_unit["gap"] < 0].sort_values("gap")
        # Cari unit dgn gap pendapatan paling minus secara keseluruhan (tanpa filter prioritas prestasi)
        target_row = gap_unit_neg.iloc[0] if not gap_unit_neg.empty else None
        if target_row is not None:
            wg = target_row
            wg_prestasi = wg["prestasi_cap"]
            prestasi_txt = f"{wg_prestasi:.1f}%" if wg_prestasi is not None else "tidak tersedia"
            if wg_prestasi is not None and wg_prestasi < 100:
                penyebab_txt = "Rendahnya capaian prestasi unit ini menjadi salah satu penyebab utama kekurangan pendapatan."
            elif wg_prestasi is not None:
                penyebab_txt = "Meski capaian prestasi sudah tercapai, gap pendapatan tetap terjadi — kemungkinan disebabkan faktor lain (tarif/rate, harga jual, atau komposisi pekerjaan)."
            else:
                penyebab_txt = "Data capaian prestasi unit ini belum tersedia untuk analisis lebih lanjut."
            add_finding_box(s, 0.6, note_top_au, 12.2, note_h_au, "⚠",
                             f"{wg['label']} adalah unit dengan GAP PENDAPATAN MINUS PALING TINGGI ({fmt_rp(wg['gap'])}) — "
                             f"Realisasi {fmt_rp(wg['pend_r'])} vs Budget {fmt_rp(wg['pend_b'])}, dengan Capaian Prestasi {prestasi_txt}. "
                             f"{penyebab_txt}",
                             RED_BG, RED, RED)
        else:
            add_finding_box(s, 0.6, note_top_au, 12.2, note_h_au, "✅",
                             "Tidak ada unit dengan gap pendapatan minus — seluruh unit mencapai/melebihi target pendapatan.",
                             GREEN_BG, GREEN, GREEN)

        # ================= SLIDE 2: AVAILABILITY — REALISASI VS BUDGET & CAPAIAN VS UTILISASI =================
        s = add_content_slide(f"AVAILABILITY — Realisasi vs Budget & Capaian s/d {period}", f"Tren Bulanan · {snum2}{divisi_label}{kat_suffix}")

        avail_rows = []
        if not sasaran_mutu_data.empty:
            sm2 = sasaran_mutu_data.copy()
            sm2 = sm2.dropna(subset=["jenis_unit"])
            avail_tbl = sm2.groupby(["lokasi", "kategori", "jenis_unit"], as_index=False).agg(
                avail_r=("availability_pct", "mean"), avail_t=("availability_target", "mean"),
                util_r=("utilisasi_pct", "mean"), util_t=("utilisasi_target", "mean"),
            )
            avail_tbl["site_short"] = avail_tbl["lokasi"].map(SITE_ABBR).fillna(avail_tbl["lokasi"])
            avail_tbl["label"] = avail_tbl["site_short"] + " — " + avail_tbl["jenis_unit"]
            avail_tbl = avail_tbl.sort_values(["lokasi", "kategori", "jenis_unit"])
            for _, r in avail_tbl.iterrows():
                avail_cap = (r["avail_r"] / r["avail_t"] * 100) if r["avail_t"] else None
                util_cap = (r["util_r"] / r["util_t"] * 100) if r["util_t"] else None
                avail_rows.append({"label": r["label"], "avail_r": r["avail_r"], "avail_t": r["avail_t"],
                                    "avail_cap": avail_cap, "util_cap": util_cap})

        panel_top2 = 1.05
        panel_bottom2 = 7.3
        total_h2 = panel_bottom2 - panel_top2
        h_chart1 = total_h2 * 0.48
        h_chart2 = total_h2 * 0.48

        def _draw_avail_chart(slide, rows, top, height, title, series1_name, series1_key, series1_color,
                               series2_name, series2_key, series2_color, num_fmt):
            add_card_panel(slide, 0.45, top, 12.35, height)
            add_panel_header(slide, 0.45, top, 12.35, title, height=0.36)
            chart_top_y = top + 0.42
            chart_h_y = height - 0.47
            if rows:
                n_y = len(rows)
                cd_y = CategoryChartData()
                cd_y.categories = [r["label"] for r in rows]
                cd_y.add_series(series1_name, tuple(round(r[series1_key], 1) if r[series1_key] is not None else 0 for r in rows))
                cd_y.add_series(series2_name, tuple(round(r[series2_key], 1) if r[series2_key] is not None else 0 for r in rows))
                gframe_y = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.6), Inches(chart_top_y), Inches(12.2), Inches(chart_h_y), cd_y)
                chart_y = gframe_y.chart
                chart_y.series[0].format.fill.solid(); chart_y.series[0].format.fill.fore_color.rgb = series1_color
                chart_y.series[1].format.fill.solid(); chart_y.series[1].format.fill.fore_color.rgb = series2_color
                chart_y.has_title = False
                plot_y = chart_y.plots[0]
                plot_y.gap_width = 60
                plot_y.has_data_labels = True
                dls_y = plot_y.data_labels
                dls_y.number_format = num_fmt; dls_y.number_format_is_linked = False
                dls_y.font.size = Pt(6.5); dls_y.font.bold = True; dls_y.font.color.rgb = TEXT_DARK; dls_y.font.name = "Calibri"
                dls_y.position = XL_LABEL_POSITION.OUTSIDE_END
                style_chart_light(chart_y, legend=True, legend_pos=XL_LEGEND_POSITION.BOTTOM)
                cat_font_y = 8.5 if n_y <= 8 else (7 if n_y <= 14 else (6 if n_y <= 22 else 5.3))
                chart_y.category_axis.tick_labels.font.size = Pt(cat_font_y)
                chart_y.value_axis.tick_labels.font.size = Pt(cat_font_y)
                chart_y.value_axis.tick_labels.number_format = num_fmt
                chart_y.value_axis.tick_labels.number_format_is_linked = False
            else:
                add_textbox(slide, 0.6, chart_top_y + 0.1, 12.0, 0.4, "Data Sasaran Mutu belum tersedia.", size=10, italic=True, color=TEXT_MUTED)

        _draw_avail_chart(s, avail_rows, panel_top2, h_chart1,
                           "🟢 Realisasi Availability vs Budget Availability — per Site & Jenis Unit",
                           "Budget Availability", "avail_t", RGBColor(0xA9, 0xB8, 0xD4),
                           "Realisasi Availability", "avail_r", TEAL, '0"%"')
        top_chart2 = panel_top2 + h_chart1 + 0.1
        _draw_avail_chart(s, avail_rows, top_chart2, h_chart2,
                           "🔵 % Capaian Availability vs % Capaian Utilisasi — per Site & Jenis Unit",
                           "% Capaian Availability", "avail_cap", TEAL,
                           "% Capaian Utilisasi", "util_cap", GOLD, '0"%"')

        # ================= SLIDE 3: BIAYA OPERASIONAL — Ringkasan Biaya vs Fisik =================
        s = add_content_slide(f"BIAYA OPERASIONAL — Budget vs Aktual s/d {period}", f"Biaya Operasional \u00b7 {snum3}{divisi_label}{kat_suffix}")

        # --- Filter khusus BBM: baris dgn qty_bbm ada TAPI biaya_bbm ATAU prestasi tidak ada -> jangan dihitung ---
        def _bbm_valid_mask(df_):
            bad_r = (df_["qty_bbm_realisasi"] > 0) & ((df_["biaya_bbm_realisasi"].fillna(0) == 0) | (df_["prestasi_realisasi"].fillna(0) == 0))
            bad_b = (df_["qty_bbm_budget"] > 0) & ((df_["biaya_bbm_budget"].fillna(0) == 0) | (df_["prestasi_budget"].fillna(0) == 0))
            return ~(bad_r | bad_b)

        data_bbm_ok = data[_bbm_valid_mask(data)]

        tot_biaya_r3 = data["total_biaya_realisasi"].sum()
        tot_biaya_b3 = data["total_biaya_budget"].sum()
        cap_biaya3 = (tot_biaya_r3 / tot_biaya_b3 * 100) if tot_biaya_b3 else None
        prestasi_r3 = data["prestasi_realisasi"].sum()
        prestasi_b3 = data["prestasi_budget"].sum()
        cap_fisik_biaya3 = (prestasi_r3 / prestasi_b3 * 100) if prestasi_b3 else None

        upah_r3 = data["upah_realisasi"].sum()
        upah_b3 = data["upah_budget"].sum()
        cap_upah3 = (upah_r3 / upah_b3 * 100) if upah_b3 else None

        bbm_biaya_r3 = data_bbm_ok["biaya_bbm_realisasi"].sum()
        bbm_biaya_b3 = data_bbm_ok["biaya_bbm_budget"].sum()
        bbm_qty_r3 = data_bbm_ok["qty_bbm_realisasi"].sum()
        bbm_qty_b3 = data_bbm_ok["qty_bbm_budget"].sum()
        # Rp/Ltr = Total Biaya BBM dibagi Total Qty BBM
        harga_bbm_r3 = (bbm_biaya_r3 / bbm_qty_r3) if bbm_qty_r3 else None
        harga_bbm_b3 = (bbm_biaya_b3 / bbm_qty_b3) if bbm_qty_b3 else None
        cap_bbm3 = (harga_bbm_r3 / harga_bbm_b3 * 100) if (harga_bbm_r3 is not None and harga_bbm_b3) else None
        bbm_prestasi_r3 = data_bbm_ok["prestasi_realisasi"].sum()
        bbm_prestasi_b3 = data_bbm_ok["prestasi_budget"].sum()

        # Capaian Fisik BBM tergantung kategori: Transportasi pakai KM/Ltr (makin tinggi makin baik),
        # Alat Berat pakai Ltr/HM (makin rendah makin baik). Kalau data campuran, tentukan dari kategori dominan.
        kategori_set_bbm3 = set(data_bbm_ok["kategori"].dropna().unique())
        if kategori_set_bbm3 == {"AB"} or (len(kategori_set_bbm3) > 1 and "ALAT BERAT" in kat_suffix):
            is_ab_bbm3 = True
        elif kategori_set_bbm3 == {"TR"} or (len(kategori_set_bbm3) > 1 and "TRANSPORTASI" in kat_suffix):
            is_ab_bbm3 = False
        elif len(kategori_set_bbm3) > 1:
            # Campuran & tidak ada penanda kategori jelas -> pakai kategori dgn kontribusi biaya BBM terbesar
            kat_biaya = data_bbm_ok.groupby("kategori")["biaya_bbm_realisasi"].sum()
            is_ab_bbm3 = (not kat_biaya.empty) and (kat_biaya.idxmax() == "AB")
        else:
            is_ab_bbm3 = False

        if is_ab_bbm3:
            # Alat Berat: Liter/HM — makin rendah makin baik
            rate_r3 = (bbm_qty_r3 / bbm_prestasi_r3) if bbm_prestasi_r3 else None
            rate_b3 = (bbm_qty_b3 / bbm_prestasi_b3) if bbm_prestasi_b3 else None
            cap_fisik_bbm3 = (rate_r3 / rate_b3 * 100) if (rate_r3 is not None and rate_b3) else None
            fisik_bbm_higher_better3 = False
        else:
            # Transportasi: KM/Ltr — makin tinggi makin baik
            rate_r3 = (bbm_prestasi_r3 / bbm_qty_r3) if bbm_qty_r3 else None
            rate_b3 = (bbm_prestasi_b3 / bbm_qty_b3) if bbm_qty_b3 else None
            cap_fisik_bbm3 = (rate_r3 / rate_b3 * 100) if (rate_r3 is not None and rate_b3) else None
            fisik_bbm_higher_better3 = True

        maint_r3 = data["maintenance_realisasi"].sum()
        maint_b3 = data["maintenance_budget"].sum()
        cap_maint3 = (maint_r3 / maint_b3 * 100) if maint_b3 else None

        lain_r3 = data["lainnya_realisasi"].sum()
        lain_b3 = data["lainnya_budget"].sum()
        cap_lain3 = (lain_r3 / lain_b3 * 100) if lain_b3 else None

        def _cap_disp3(v):
            if v is None or pd.isna(v):
                return "-"
            ok = "\u2713" if v <= 100 else "\u2717"
            vv = ">999%" if v > 999 else f"{v:.1f}%"
            return f"{ok} {vv}"

        def _fisik_disp3(v, higher_is_better):
            if v is None or pd.isna(v):
                return "-"
            good = (v >= 100) if higher_is_better else (v <= 100)
            ok = "\u2713" if good else "\u2717"
            vv = ">999%" if v > 999 else f"{v:.1f}%"
            return f"{ok} {vv}"

        ringkasan3_rows = [
            ["Total Biaya", fmt_rp(tot_biaya_b3), fmt_rp(tot_biaya_r3), _cap_disp3(cap_biaya3), _fisik_disp3(cap_fisik_biaya3, True)],
            ["Upah Operator", fmt_rp(upah_b3), fmt_rp(upah_r3), _cap_disp3(cap_upah3), "-"],
            ["Biaya BBM (Rp/Ltr)",
             f"Rp {harga_bbm_b3:,.0f}" if harga_bbm_b3 is not None else "-",
             f"Rp {harga_bbm_r3:,.0f}" if harga_bbm_r3 is not None else "-",
             _cap_disp3(cap_bbm3), "-"],
            ["Biaya Maintenance", fmt_rp(maint_b3), fmt_rp(maint_r3), _cap_disp3(cap_maint3), "-"],
            ["Biaya Lainnya", fmt_rp(lain_b3), fmt_rp(lain_r3), _cap_disp3(cap_lain3), "-"],
        ]

        add_textbox(s, 0.4, 0.98, 5.9, 0.3, f"Ringkasan Biaya PT. BKMS (s/d {period})", size=14, bold=True, color=TEXT_DARK)

        tbl3_top = 1.33
        tbl3_h = 2.75
        add_table(s, 0.4, tbl3_top, 5.9, tbl3_h,
                  ["Metrik", "Budget", "Aktual", "Capaian", "Cap. Fisik"], ringkasan3_rows,
                  status_col=[3, 4], col_widths=[1.7, 1.15, 1.15, 0.95, 0.95], font_size=10, header_size=10,
                  fill_badge=True)

        # --- Catatan otomatis: metrik biaya mana yang paling over budget ---
        cap_map3 = {"Total Biaya": cap_biaya3, "Upah Operator": cap_upah3, "Biaya BBM (Rp/Ltr)": cap_bbm3,
                    "Biaya Maintenance": cap_maint3, "Biaya Lainnya": cap_lain3}
        over_items3 = {k: v for k, v in cap_map3.items() if v is not None and v > 100 and k != "Total Biaya"}
        note_top3 = tbl3_top + tbl3_h + 0.15
        note_h3 = 0.75
        if over_items3:
            worst_label3 = max(over_items3, key=over_items3.get)
            worst_val3 = over_items3[worst_label3]
            add_finding_box(s, 0.4, note_top3, 5.9, note_h3, "\u26A0",
                             f"{worst_label3} OVER BUDGET ({worst_val3:.1f}%) \u2014 perlu efisiensi biaya s/d {period}.",
                             RED_BG, RED, RED)
        else:
            add_finding_box(s, 0.4, note_top3, 5.9, note_h3, "\u2705",
                             "Seluruh komponen biaya berada dalam/di bawah budget.",
                             GREEN_BG, GREEN, GREEN)
        left_col_bottom3 = note_top3 + note_h3

        # ================= PANEL KANAN ATAS: BTL per Site & Kategori =================
        btl_sk3 = data.groupby(["lokasi", "kategori"], as_index=False).agg(
            btl_r=("biaya_tidak_langsung_realisasi", "sum"), btl_b=("biaya_tidak_langsung_budget", "sum"))
        btl_sk3 = btl_sk3[(btl_sk3["btl_r"] > 0) | (btl_sk3["btl_b"] > 0)].copy()
        total_btl_r3 = btl_sk3["btl_r"].sum()
        btl_sk3["site_short"] = btl_sk3["lokasi"].map(SITE_ABBR).fillna(btl_sk3["lokasi"])
        btl_sk3["label"] = btl_sk3["site_short"] + " (" + btl_sk3["kategori"] + ")"
        btl_sk3 = btl_sk3.sort_values("btl_r", ascending=False)

        over_btl_sites3 = []
        btl3_rows = []
        for _, r in btl_sk3.iterrows():
            pct_target = (r["btl_r"] / r["btl_b"] * 100) if r["btl_b"] else None
            pct_share = (r["btl_r"] / total_btl_r3 * 100) if total_btl_r3 else 0
            pct_disp = (f"\u2713 {pct_target:.0f}%" if pct_target is not None and pct_target <= 100
                        else (f"\u2717 {pct_target:.0f}%" if pct_target is not None else "-"))
            btl3_rows.append([r["label"], fmt_rp(r["btl_b"]), fmt_rp(r["btl_r"]), pct_disp, f"{pct_share:.1f}%"])
            if pct_target is not None and pct_target > 100:
                over_btl_sites3.append(r["label"])

        btl_panel_top3 = 0.98
        n_btl3 = max(len(btl3_rows), 1)
        row_h_btl3 = 0.32 if n_btl3 <= 6 else (0.26 if n_btl3 <= 10 else 0.22)
        tbl_h_btl3 = row_h_btl3 * (n_btl3 + 1)
        btl_content_h3 = 0.15 + 0.45 + 0.1 + tbl_h_btl3 + 0.15
        # Regangkan tinggi card BTL supaya sejajar dgn panel kiri (Ringkasan Biaya), tidak menyisakan
        # celah kosong sebelum panel Maintenance di bawahnya. Konten (banner+tabel) tetap di posisi natural.
        btl_panel_h3 = max(btl_content_h3, left_col_bottom3 - btl_panel_top3)
        add_card_panel(s, 6.85, btl_panel_top3, 6.05, btl_panel_h3)
        if not over_btl_sites3:
            add_status_banner(s, 7.1, btl_panel_top3 + 0.15, 5.55, 0.45, "\u2705", "BTL \u2014 UNDER BUDGET secara keseluruhan", GREEN_BG, GREEN, GREEN)
        else:
            over_txt3 = " & ".join(over_btl_sites3[:3]) + (", dll" if len(over_btl_sites3) > 3 else "")
            add_status_banner(s, 7.1, btl_panel_top3 + 0.15, 5.55, 0.45, "\u26a0\ufe0f", f"BTL \u2014 UNDER BUDGET, kecuali {over_txt3}", RED_BG, RED, RED)
        font_btl3 = 9.5 if n_btl3 <= 6 else (8.5 if n_btl3 <= 10 else 7.5)
        add_table(s, 7.1, btl_panel_top3 + 0.7, 5.55, tbl_h_btl3, ["Site (Kategori)", "Budget", "Aktual", "% Target", "% BTL"], btl3_rows,
                  status_col=3, col_widths=[1.55, 1.15, 1.15, 0.85, 0.85], font_size=font_btl3, header_size=font_btl3)
        right_col_bottom3 = btl_panel_top3 + btl_panel_h3

        # ================= PANEL BAWAH (LEBAR PENUH): % Capaian Maintenance per Site & Jenis Unit =================
        maint_su3 = data.groupby(["lokasi", "kategori", "jenis_unit"], as_index=False).agg(
            maint_r=("maintenance_realisasi", "sum"), maint_b=("maintenance_budget", "sum"))
        maint_su3 = maint_su3[maint_su3["maint_b"] > 0].copy()
        maint_su3["site_short"] = maint_su3["lokasi"].map(SITE_ABBR).fillna(maint_su3["lokasi"])
        maint_su3["label"] = maint_su3["site_short"] + " \u2014 " + maint_su3["jenis_unit"]
        maint_su3["cap"] = maint_su3["maint_r"] / maint_su3["maint_b"] * 100
        maint_su3["gap"] = maint_su3["maint_r"] - maint_su3["maint_b"]
        maint_su3 = maint_su3.sort_values("gap", ascending=False)
        n_maint3 = max(len(maint_su3), 1)

        maint_panel_top3 = max(left_col_bottom3, right_col_bottom3) + 0.15
        maint_panel_h3 = 7.3 - maint_panel_top3
        add_card_panel(s, 0.4, maint_panel_top3, 12.5, maint_panel_h3)
        add_panel_header(s, 0.4, maint_panel_top3, 12.5, "\u2699 % Capaian Maintenance \u2014 per Site & Jenis Unit", height=0.36)
        chart_top_m3 = maint_panel_top3 + 0.42
        chart_h_m3 = maint_panel_h3 - 0.47
        if not maint_su3.empty:
            cd_m3 = CategoryChartData()
            cd_m3.categories = list(maint_su3["label"])
            cd_m3.add_series("% Capaian Maintenance", tuple(round(v, 1) for v in maint_su3["cap"]))
            gframe_m3 = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.55), Inches(chart_top_m3), Inches(12.2), Inches(chart_h_m3), cd_m3)
            chart_m3 = gframe_m3.chart
            chart_m3.series[0].format.fill.solid(); chart_m3.series[0].format.fill.fore_color.rgb = TEAL
            chart_m3.has_title = False
            plot_m3 = chart_m3.plots[0]
            plot_m3.gap_width = 50
            label_font_m3 = 7.5 if n_maint3 <= 12 else (6.5 if n_maint3 <= 20 else 6)
            from pptx.oxml.ns import qn as _qn
            for i, pt in enumerate(chart_m3.series[0].points):
                v = maint_su3["cap"].iloc[i]
                if v > 105:
                    pt.format.fill.solid(); pt.format.fill.fore_color.rgb = RED
                r_val = maint_su3["maint_r"].iloc[i]
                b_val = maint_su3["maint_b"].iloc[i]
                gap_val = r_val - b_val
                gap_sign = "+" if gap_val >= 0 else "-"
                dl = pt.data_label
                dl.has_text_frame = True
                if n_maint3 <= 12:
                    dl.text_frame.text = f"{v:.0f}% ({gap_sign}{fmt_rp(abs(gap_val))})"
                else:
                    dl.text_frame.text = f"{v:.0f}% ({gap_sign}{fmt_rp(abs(gap_val))})"
                    # Kategori banyak: putar teks label vertikal (90 derajat) supaya tidak numpuk horizontal
                    bodyPr = dl.text_frame._txBody.find(_qn('a:bodyPr'))
                    if bodyPr is not None:
                        bodyPr.set('rot', '-5400000')
                        bodyPr.set('vert', 'horz')
                r0 = dl.text_frame.paragraphs[0].runs[0]
                r0.font.size = Pt(label_font_m3); r0.font.bold = True; r0.font.color.rgb = TEXT_DARK; r0.font.name = "Calibri"
            style_chart_light(chart_m3, legend=False)
            cat_font_m3 = 8 if n_maint3 <= 10 else (6.5 if n_maint3 <= 20 else 5.3)
            chart_m3.category_axis.tick_labels.font.size = Pt(cat_font_m3)
            chart_m3.value_axis.tick_labels.font.size = Pt(cat_font_m3)
        else:
            add_textbox(s, 0.55, chart_top_m3 + 0.1, 12.0, 0.5, "Data Maintenance belum tersedia.", size=10, italic=True, color=TEXT_MUTED)


        # ================= SLIDE 4: ANALISIS Capaian Biaya vs Capaian Fisik & Konsumsi BBM =================
        s = add_content_slide(f"ANALISIS: Capaian Biaya vs Capaian Fisik \u2014 s/d {period}", f"Analisis Biaya \u00b7 {snum4}{divisi_label}{kat_suffix}")

        panel_top4 = 1.0
        panel_bottom4 = 7.3
        panel_h4 = panel_bottom4 - panel_top4

        # ================= PANEL KIRI: Capaian Biaya (di luar Penyusutan) vs Capaian Prestasi =================
        prestasi_r4 = data["prestasi_realisasi"].sum()
        prestasi_b4 = data["prestasi_budget"].sum()
        cap_prestasi_global4 = (prestasi_r4 / prestasi_b4 * 100) if prestasi_b4 else None

        comp_defs4 = [
            ("Upah Operator", "upah_realisasi", "upah_budget"),
            ("Biaya BBM", "biaya_bbm_realisasi", "biaya_bbm_budget"),
            ("Biaya Maintenance", "maintenance_realisasi", "maintenance_budget"),
            ("Biaya Lainnya", "lainnya_realisasi", "lainnya_budget"),
        ]
        left_rows4 = []
        for name4, rcol4, bcol4 in comp_defs4:
            r4 = data[rcol4].sum()
            b4 = data[bcol4].sum()
            cap_biaya4 = (r4 / b4 * 100) if b4 else None
            left_rows4.append({"label": name4, "cap_biaya": cap_biaya4, "cap_prestasi": cap_prestasi_global4})
        left_rows4 = sorted(left_rows4, key=lambda r: (r["cap_biaya"] if r["cap_biaya"] is not None else -1), reverse=True)

        add_card_panel(s, 0.4, panel_top4, 6.05, panel_h4)
        add_panel_header(s, 0.4, panel_top4, 6.05, "$ Capaian Biaya (di luar Penyusutan) vs Capaian Prestasi", height=0.4)
        chart_top_l4 = panel_top4 + 0.45
        chart_h_l4 = panel_h4 - 0.55
        cd_l4 = CategoryChartData()
        cd_l4.categories = [r["label"] for r in left_rows4]
        cd_l4.add_series("% Capaian Biaya", tuple(round(r["cap_biaya"], 1) if r["cap_biaya"] is not None else 0 for r in left_rows4))
        cd_l4.add_series("% Capaian Prestasi", tuple(round(r["cap_prestasi"], 1) if r["cap_prestasi"] is not None else 0 for r in left_rows4))
        gframe_l4 = s.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, Inches(0.55), Inches(chart_top_l4), Inches(5.75), Inches(chart_h_l4), cd_l4)
        chart_l4 = gframe_l4.chart
        chart_l4.series[0].format.fill.solid(); chart_l4.series[0].format.fill.fore_color.rgb = RED
        chart_l4.series[1].format.fill.solid(); chart_l4.series[1].format.fill.fore_color.rgb = TEAL
        chart_l4.has_title = False
        plot_l4 = chart_l4.plots[0]
        plot_l4.has_data_labels = True
        dls_l4 = plot_l4.data_labels
        dls_l4.number_format = '0.0"%"'; dls_l4.number_format_is_linked = False
        dls_l4.font.size = Pt(9.5); dls_l4.font.bold = True; dls_l4.font.color.rgb = TEXT_DARK; dls_l4.font.name = "Calibri"
        style_chart_light(chart_l4, legend=True, legend_pos=XL_LEGEND_POSITION.TOP)
        chart_l4.category_axis.tick_labels.font.size = Pt(10.5)
        chart_l4.value_axis.tick_labels.font.size = Pt(9.5)

        # ================= PANEL KANAN: % Capaian Konsumsi BBM per Site & Jenis Unit =================
        # Filter (per baris, realisasi & budget dicek terpisah):
        #  - qty BBM ada tapi Rp BBM tidak ada -> qty & prestasi baris tsb tidak ikut dihitung
        #  - prestasi ada tapi qty BBM tidak ada -> prestasi & qty baris tsb tidak ikut dihitung
        #  - qty BBM ada tapi prestasi tidak ada -> qty & prestasi baris tsb tidak ikut dihitung
        data_bbm_ok4 = data.copy()

        valid_r4 = (data_bbm_ok4["qty_bbm_realisasi"].fillna(0) > 0) & \
                   (data_bbm_ok4["biaya_bbm_realisasi"].fillna(0) > 0) & \
                   (data_bbm_ok4["prestasi_realisasi"].fillna(0) > 0)
        partial_r4 = (~valid_r4) & ((data_bbm_ok4["qty_bbm_realisasi"].fillna(0) > 0) | (data_bbm_ok4["prestasi_realisasi"].fillna(0) > 0))
        data_bbm_ok4.loc[partial_r4, ["qty_bbm_realisasi", "prestasi_realisasi"]] = 0

        valid_b4 = (data_bbm_ok4["qty_bbm_budget"].fillna(0) > 0) & \
                   (data_bbm_ok4["biaya_bbm_budget"].fillna(0) > 0) & \
                   (data_bbm_ok4["prestasi_budget"].fillna(0) > 0)
        partial_b4 = (~valid_b4) & ((data_bbm_ok4["qty_bbm_budget"].fillna(0) > 0) | (data_bbm_ok4["prestasi_budget"].fillna(0) > 0))
        data_bbm_ok4.loc[partial_b4, ["qty_bbm_budget", "prestasi_budget"]] = 0

        bbm_su4 = data_bbm_ok4.groupby(["lokasi", "kategori", "jenis_unit"], as_index=False).agg(
            qty_r=("qty_bbm_realisasi", "sum"), qty_b=("qty_bbm_budget", "sum"),
            prestasi_r=("prestasi_realisasi", "sum"), prestasi_b=("prestasi_budget", "sum"))
        bbm_su4 = bbm_su4[(bbm_su4["qty_r"] > 0) | (bbm_su4["qty_b"] > 0)].copy()
        bbm_su4["site_short"] = bbm_su4["lokasi"].map(SITE_ABBR).fillna(bbm_su4["lokasi"])
        bbm_su4["label"] = bbm_su4["site_short"] + " \u2014 " + bbm_su4["jenis_unit"]

        # Biaya BBM (Rp) per site & jenis unit -> dipakai utk urutan (Rupiah over tertinggi ditampilkan pertama)
        biaya_bbm_su4 = data.groupby(["lokasi", "jenis_unit"], as_index=False).agg(
            biaya_r=("biaya_bbm_realisasi", "sum"), biaya_b=("biaya_bbm_budget", "sum"))
        biaya_bbm_su4["gap_rp"] = biaya_bbm_su4["biaya_r"] - biaya_bbm_su4["biaya_b"]
        gap_rp_lookup4 = {(r["lokasi"], r["jenis_unit"]): r["gap_rp"] for _, r in biaya_bbm_su4.iterrows()}
        bbm_su4["gap_rp"] = bbm_su4.apply(lambda r: gap_rp_lookup4.get((r["lokasi"], r["jenis_unit"]), 0), axis=1)

        def _bbm_cap4(row):
            if row["kategori"] == "AB":
                rate_r = (row["qty_r"] / row["prestasi_r"]) if row["prestasi_r"] else None
                rate_b = (row["qty_b"] / row["prestasi_b"]) if row["prestasi_b"] else None
            else:
                rate_r = (row["prestasi_r"] / row["qty_r"]) if row["qty_r"] else None
                rate_b = (row["prestasi_b"] / row["qty_b"]) if row["qty_b"] else None
            if rate_r is None or not rate_b:
                return None
            return rate_r / rate_b * 100

        bbm_su4["cap"] = bbm_su4.apply(_bbm_cap4, axis=1)
        bbm_su4 = bbm_su4.dropna(subset=["cap"])
        bbm_su4 = bbm_su4.sort_values("gap_rp", ascending=False)
        n_bbm4 = max(len(bbm_su4), 1)

        add_card_panel(s, 6.85, panel_top4, 6.05, panel_h4)
        add_panel_header(s, 6.85, panel_top4, 6.05, "\u26fd % Capaian Konsumsi BBM \u2014 per Site & Jenis Unit", height=0.4)
        chart_top_r4 = panel_top4 + 0.45
        note_h_r4 = 1.05
        chart_h_r4 = panel_h4 - 0.45 - 0.15 - note_h_r4 - 0.1
        if not bbm_su4.empty:
            cd_r4 = CategoryChartData()
            cd_r4.categories = list(bbm_su4["label"])
            cd_r4.add_series("% Capaian Konsumsi BBM", tuple(round(v, 1) for v in bbm_su4["cap"]))
            gframe_r4 = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(7.0), Inches(chart_top_r4), Inches(5.75), Inches(chart_h_r4), cd_r4)
            chart_r4 = gframe_r4.chart
            chart_r4.series[0].format.fill.solid(); chart_r4.series[0].format.fill.fore_color.rgb = TEAL
            chart_r4.has_title = False
            plot_r4 = chart_r4.plots[0]
            plot_r4.gap_width = 50
            label_font_r4 = 8 if n_bbm4 <= 6 else (7.5 if n_bbm4 <= 12 else 6)
            from pptx.oxml.ns import qn as _qn4
            avg_abs_gap_r4 = bbm_su4["gap_rp"].abs().mean() if not bbm_su4.empty else 0
            for i, pt in enumerate(chart_r4.series[0].points):
                v = bbm_su4["cap"].iloc[i]
                gap_val4 = bbm_su4["gap_rp"].iloc[i]
                if avg_abs_gap_r4 and gap_val4 > 0 and gap_val4 > avg_abs_gap_r4:
                    pt.format.fill.solid(); pt.format.fill.fore_color.rgb = RED
                gap_sign4 = "+" if gap_val4 >= 0 else "-"
                dl = pt.data_label
                dl.has_text_frame = True
                tf = dl.text_frame
                tf.text = f"{v:.0f}%"
                p2 = tf.add_paragraph()
                p2.text = f"({gap_sign4}{fmt_rp(abs(gap_val4))})"
                for para in tf.paragraphs:
                    for run in para.runs:
                        run.font.size = Pt(label_font_r4); run.font.bold = True; run.font.color.rgb = TEXT_DARK; run.font.name = "Calibri"
                if n_bbm4 > 6:
                    bodyPr = dl.text_frame._txBody.find(_qn4('a:bodyPr'))
                    if bodyPr is not None:
                        bodyPr.set('rot', '-5400000')
            style_chart_light(chart_r4, legend=False)
            cat_font_r4 = 8 if n_bbm4 <= 10 else (6.5 if n_bbm4 <= 20 else 5.3)
            chart_r4.category_axis.tick_labels.font.size = Pt(cat_font_r4)
            chart_r4.value_axis.tick_labels.font.size = Pt(cat_font_r4)

            # --- Analisa: kenapa unit #1 (gap Rupiah terbesar) bisa dominan, meski %capaian konsumsinya kecil ---
            top4 = bbm_su4.iloc[0]
            unit_rows4 = data[(data["lokasi"] == top4["lokasi"]) & (data["jenis_unit"] == top4["jenis_unit"])]
            biaya_r_u4 = unit_rows4["biaya_bbm_realisasi"].sum()
            biaya_b_u4 = unit_rows4["biaya_bbm_budget"].sum()
            qty_r_u4 = unit_rows4["qty_bbm_realisasi"].sum()
            qty_b_u4 = unit_rows4["qty_bbm_budget"].sum()
            harga_r_u4 = (biaya_r_u4 / qty_r_u4) if qty_r_u4 else None
            harga_b_u4 = (biaya_b_u4 / qty_b_u4) if qty_b_u4 else None
            harga_pct_u4 = (harga_r_u4 / harga_b_u4 * 100) if (harga_r_u4 is not None and harga_b_u4) else None
            note_top_r4 = chart_top_r4 + chart_h_r4 + 0.15
            if harga_pct_u4 is not None:
                harga_selisih4 = harga_pct_u4 - 100
                if abs(harga_selisih4) > abs(top4["cap"] - 100):
                    sebab_txt4 = (f"harga BBM realisasi (Rp {harga_r_u4:,.0f}/Ltr) yang {'lebih tinggi' if harga_selisih4 > 0 else 'lebih rendah'} "
                                  f"{abs(harga_selisih4):.1f}% dari budget (Rp {harga_b_u4:,.0f}/Ltr) \u2014 bukan konsumsinya.")
                else:
                    sebab_txt4 = (f"volume konsumsi BBM yang besar ({qty_r_u4:,.0f} Ltr), sehingga meski capaian konsumsi hanya "
                                  f"{top4['cap']:.1f}%, dampak Rupiah-nya tetap signifikan.")
                add_finding_box(s, 7.0, note_top_r4, 5.6, note_h_r4, "\u2731",
                                 f"{top4['label']} adalah dampak Rupiah terbesar ({'+' if top4['gap_rp'] >= 0 else '-'}{fmt_rp(abs(top4['gap_rp']))}), "
                                 f"walau capaian konsumsi hanya {top4['cap']:.1f}% (dari target 100%). Penyebab utamanya adalah {sebab_txt4}",
                                 GOLD_BG, GOLD, RGBColor(0x7A, 0x5C, 0x0D))
        else:
            add_textbox(s, 7.0, chart_top_r4 + 0.1, 5.6, 0.5, "Data konsumsi BBM belum tersedia.", size=10, italic=True, color=TEXT_MUTED)

        # ================= SLIDE 5: KEY INSIGHTS \u2014 DOWNTIME ANALYSIS & VARIAN =================
        s = add_content_slide(f"KEY INSIGHTS \u2014 Downtime Analysis & Varian s/d {period}", f"Analisis Downtime \u00b7 {snum5}{divisi_label}{kat_suffix}")

        dt_avg_r5 = sasaran_mutu_data["downtime_pct"].mean() if not sasaran_mutu_data.empty else None
        dt_avg_t5 = sasaran_mutu_data["downtime_target"].mean() if not sasaran_mutu_data.empty else None
        cap_dt5 = (dt_avg_r5 / dt_avg_t5 * 100) if (dt_avg_r5 is not None and dt_avg_t5) else None
        varian_dt5 = (dt_avg_r5 - dt_avg_t5) if (dt_avg_r5 is not None and dt_avg_t5 is not None) else None
        good_dt5 = varian_dt5 is not None and varian_dt5 <= 0
        avail_target5 = (100 - dt_avg_t5) if dt_avg_t5 is not None else None
        avail_aktual5 = (100 - dt_avg_r5) if dt_avg_r5 is not None else None

        # ================= BARIS ATAS: 3 KARTU KPI =================
        card_top5 = 0.98
        card_h5 = 1.85
        card_gap5 = 0.25
        card_w5 = (12.5 - 2 * card_gap5) / 3

        add_kpi_card(s, 0.4, card_top5, card_w5, card_h5, "\u23f8", RED, RED,
                     "Downtime Aktual (s/d " + period + ")",
                     (f"{dt_avg_r5:.2f}%" if dt_avg_r5 is not None else "-"),
                     (f"Target: {dt_avg_t5:.2f}%  |  Varian: {varian_dt5:+.2f}%" if varian_dt5 is not None else "-"),
                     (f"\u2717 VARIAN {varian_dt5:.2f}% BELUM TERKENDALI" if (varian_dt5 is not None and varian_dt5 > 0)
                      else (f"\u2713 DALAM TARGET" if varian_dt5 is not None else "Data tidak tersedia")),
                     good_dt5)

        add_kpi_card(s, 0.4 + card_w5 + card_gap5, card_top5, card_w5, card_h5, "\u2699", TEAL, TEAL,
                     "Target Downtime (Diizinkan)",
                     (f"{dt_avg_t5:.2f}%" if dt_avg_t5 is not None else "-"),
                     (f"Availability Target: {avail_target5:.2f}%" if avail_target5 is not None else "-"),
                     "BATAS MAKSIMUM DOWNTIME",
                     True)

        cap_good5 = cap_dt5 is not None and cap_dt5 <= 100
        add_kpi_card(s, 0.4 + 2 * (card_w5 + card_gap5), card_top5, card_w5, card_h5, "\u25CE", GREEN if cap_good5 else RED, GREEN if cap_good5 else RED,
                     "% Capaian Realisasi Downtime",
                     (f"{cap_dt5:.1f}%" if cap_dt5 is not None else "-"),
                     (f"Availability Aktual: {avail_aktual5:.2f}%" if avail_aktual5 is not None else "-"),
                     (f"{cap_dt5:.1f}% \u2014 {'Under' if cap_good5 else 'Over'} Budget" if cap_dt5 is not None else "Data tidak tersedia"),
                     cap_good5)

        # ================= BARIS BAWAH =================
        panel_top5 = card_top5 + card_h5 + 0.15
        panel_bottom5 = 7.3
        panel_h5 = panel_bottom5 - panel_top5
        left_w5 = 7.3
        right_x5 = 7.9
        right_w5 = 5.0

        # --- Panel kiri: % Downtime per Site & Jenis Unit + catatan strategi ---
        add_card_panel(s, 0.4, panel_top5, left_w5, panel_h5, accent_color=RED)
        add_panel_header(s, 0.4, panel_top5, left_w5, "\u23f8 % Downtime \u2014 per Site & Jenis Unit (Min Target: " +
                          (f"{dt_avg_t5:.2f}%" if dt_avg_t5 is not None else "-") + ")", height=0.4)
        chart_top5 = panel_top5 + 0.45
        note_h5 = 0.85
        chart_h5 = panel_h5 - 0.45 - note_h5 - 0.15
        dt_su5 = pd.DataFrame()
        if not sasaran_mutu_data.empty:
            dt_su5 = sasaran_mutu_data.dropna(subset=["jenis_unit"]).groupby(["lokasi", "kategori", "jenis_unit"], as_index=False).agg(
                dt_r=("downtime_pct", "mean"), dt_t=("downtime_target", "mean"))
            dt_su5["site_short"] = dt_su5["lokasi"].map(SITE_ABBR).fillna(dt_su5["lokasi"])
            dt_su5["label"] = dt_su5["site_short"] + " \u2014 " + dt_su5["jenis_unit"]
            dt_su5 = dt_su5.sort_values("dt_r", ascending=False)
            n_dt5 = max(len(dt_su5), 1)

            cd_dt5 = CategoryChartData()
            cd_dt5.categories = list(dt_su5["label"])
            cd_dt5.add_series("Realisasi Downtime (%)", tuple(round(v, 2) for v in dt_su5["dt_r"]))
            cd_dt5.add_series("Target Downtime (%)", tuple(round(v, 2) for v in dt_su5["dt_t"]))
            gframe_dt5 = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.55), Inches(chart_top5), Inches(7.0), Inches(chart_h5), cd_dt5)
            chart_dt5 = gframe_dt5.chart
            chart_dt5.series[0].format.fill.solid(); chart_dt5.series[0].format.fill.fore_color.rgb = RED
            chart_dt5.series[1].format.fill.solid(); chart_dt5.series[1].format.fill.fore_color.rgb = RGBColor(0xA9, 0xB8, 0xD4)
            chart_dt5.has_title = False
            for i, pt in enumerate(chart_dt5.series[0].points):
                dt_r_v = dt_su5["dt_r"].iloc[i]
                dt_t_v = dt_su5["dt_t"].iloc[i]
                if dt_r_v <= dt_t_v:
                    pt.format.fill.solid(); pt.format.fill.fore_color.rgb = GREEN
            plot_dt5 = chart_dt5.plots[0]
            plot_dt5.gap_width = 50
            plot_dt5.has_data_labels = True
            dls_dt5 = plot_dt5.data_labels
            dls_dt5.number_format = '0.0"%"'; dls_dt5.number_format_is_linked = False
            label_font_dt5 = 7 if n_dt5 <= 12 else (5.5 if n_dt5 <= 25 else 4.3)
            dls_dt5.font.size = Pt(label_font_dt5); dls_dt5.font.bold = True; dls_dt5.font.color.rgb = TEXT_DARK; dls_dt5.font.name = "Calibri"
            dls_dt5.position = XL_LABEL_POSITION.OUTSIDE_END
            style_chart_light(chart_dt5, legend=True, legend_pos=XL_LEGEND_POSITION.TOP)
            cat_font_dt5 = 7.5 if n_dt5 <= 10 else (6 if n_dt5 <= 20 else (5 if n_dt5 <= 30 else 4.2))
            chart_dt5.category_axis.tick_labels.font.size = Pt(cat_font_dt5)
            chart_dt5.value_axis.tick_labels.font.size = Pt(cat_font_dt5)
            if n_dt5 > 20:
                plot_dt5.gap_width = 30
        else:
            add_textbox(s, 0.55, chart_top5 + 0.1, 6.9, 0.5, "Data Downtime belum tersedia.", size=10, italic=True, color=TEXT_MUTED)

        # --- Catatan strategi: unit paling kritis (gap downtime terbesar) ---
        note_top5 = chart_top5 + chart_h5 + 0.1
        if not dt_su5.empty:
            worst5 = dt_su5.iloc[0]
            gap_worst5 = worst5["dt_r"] - worst5["dt_t"]
            add_finding_box(s, 0.55, note_top5, 6.9, note_h5, "\u26A0",
                             f"Strategi Perbaikan: Turunkan Downtime {worst5['label']} dari {worst5['dt_r']:.1f}% \u2192 "
                             f"\u2264{worst5['dt_t']:.1f}% (satu-satunya/unit paling kritis OVER target, gap {gap_worst5:+.1f}%). "
                             f"Cek riwayat kerusakan & jadwal preventive maintenance unit ini.",
                             RED_BG, RED, RED)
        else:
            add_finding_box(s, 0.55, note_top5, 6.9, note_h5, "\u26A0",
                             "Data unit belum tersedia untuk rekomendasi strategi perbaikan.",
                             GOLD_BG, GOLD, RGBColor(0x7A, 0x5C, 0x0D))

        # --- Panel kanan: Downtime Chain (dari Target ke Realisasi) + Strategi ---
        add_card_panel(s, right_x5, panel_top5, right_w5, panel_h5, accent_color=GOLD)
        add_panel_header(s, right_x5, panel_top5, right_w5, "\u26D3 Downtime Chain \u2014 Dari Target ke Realisasi", height=0.4)

        chain_top5 = panel_top5 + 0.55
        chain_row_h5 = 0.62
        chain_rows5 = [
            ("Target Downtime", (f"{dt_avg_t5:.2f}%" if dt_avg_t5 is not None else "-"), RGBColor(0xDD, 0xE7, 0xF7), TEAL,
             "Batas maksimum yg diizinkan"),
            ("(+) Downtime Berlebih", (f"{max(0, varian_dt5):.2f}%" if varian_dt5 is not None else "-"), RGBColor(0xFD, 0xE9, 0xD9), GOLD,
             "Kelebihan dari target"),
            ("= Realisasi Downtime", (f"{dt_avg_r5:.2f}%" if dt_avg_r5 is not None else "-"), RGBColor(0xE8, 0xE3, 0xF7), RGBColor(0x7B, 0x5C, 0xE8),
             "Yang benar-benar terjadi"),
            ("Varian (gap)", (f"{varian_dt5:+.2f}%" if varian_dt5 is not None else "-"), RED_BG, RED,
             f"Availability turun {(avail_target5 - avail_aktual5):.2f}%" if (avail_target5 is not None and avail_aktual5 is not None) else "-"),
        ]
        for i5, (lbl5, val5, bg5, bord5, note5) in enumerate(chain_rows5):
            ry5c = chain_top5 + i5 * (chain_row_h5 + 0.08)
            box5 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(right_x5 + 0.15), Inches(ry5c), Inches(2.1), Inches(chain_row_h5))
            box5.adjustments[0] = 0.12
            box5.fill.solid(); box5.fill.fore_color.rgb = bg5
            box5.line.color.rgb = bord5; box5.line.width = Pt(1)
            box5.shadow.inherit = False
            btf5 = box5.text_frame; btf5.word_wrap = True; btf5.vertical_anchor = MSO_ANCHOR.MIDDLE
            btf5.margin_left = Inches(0.1); btf5.margin_right = Inches(0.05)
            bp5 = btf5.paragraphs[0]
            br5 = bp5.add_run(); br5.text = lbl5
            br5.font.size = Pt(9); br5.font.bold = True; br5.font.color.rgb = bord5; br5.font.name = "Calibri"

            val_box5 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(right_x5 + 2.35), Inches(ry5c), Inches(1.15), Inches(chain_row_h5))
            val_box5.adjustments[0] = 0.12
            val_box5.fill.solid(); val_box5.fill.fore_color.rgb = bg5
            val_box5.line.color.rgb = bord5; val_box5.line.width = Pt(1)
            val_box5.shadow.inherit = False
            vtf5 = val_box5.text_frame; vtf5.vertical_anchor = MSO_ANCHOR.MIDDLE
            vp5 = vtf5.paragraphs[0]; vp5.alignment = PP_ALIGN.CENTER
            vr5 = vp5.add_run(); vr5.text = val5
            vr5.font.size = Pt(13); vr5.font.bold = True; vr5.font.color.rgb = bord5; vr5.font.name = "Calibri"

            add_textbox(s, right_x5 + 3.6, ry5c, right_w5 - 3.75, chain_row_h5, note5, size=7.5, color=TEXT_MUTED)

        # --- Strategi Jangka Pendek & Panjang ---
        strat_top5 = chain_top5 + 4 * (chain_row_h5 + 0.08) + 0.05
        strat_h5 = panel_h5 - (strat_top5 - panel_top5) - 0.1
        strat_box5 = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(right_x5 + 0.15), Inches(strat_top5), Inches(right_w5 - 0.3), Inches(max(strat_h5, 0.9)))
        strat_box5.fill.solid(); strat_box5.fill.fore_color.rgb = RED_BG
        strat_box5.line.color.rgb = RED; strat_box5.line.width = Pt(1)
        strat_box5.shadow.inherit = False
        stf5 = strat_box5.text_frame; stf5.word_wrap = True; stf5.margin_left = Inches(0.12); stf5.margin_top = Inches(0.08)
        sp1_5 = stf5.paragraphs[0]
        sr1_5 = sp1_5.add_run(); sr1_5.text = "\u26A0 Strategi Jangka Pendek:"
        sr1_5.font.size = Pt(9); sr1_5.font.bold = True; sr1_5.font.color.rgb = RED; sr1_5.font.name = "Calibri"
        sp2_5 = stf5.add_paragraph()
        sr2_5 = sp2_5.add_run()
        sr2_5.text = (f"Fokus perbaikan {dt_su5.iloc[0]['label']}, percepat preventive maintenance & "
                       f"minimalkan unschedule downtime, target eliminasi varian {varian_dt5:.2f}%." if not dt_su5.empty and varian_dt5 is not None
                       else "Percepat preventive maintenance & minimalkan unschedule downtime.")
        sr2_5.font.size = Pt(8.5); sr2_5.font.color.rgb = TEXT_DARK; sr2_5.font.name = "Calibri"
        sp3_5 = stf5.add_paragraph(); sp3_5.space_before = Pt(6)
        sr3_5 = sp3_5.add_run(); sr3_5.text = "\u25C6 Strategi Jangka Panjang:"
        sr3_5.font.size = Pt(9); sr3_5.font.bold = True; sr3_5.font.color.rgb = RED; sr3_5.font.name = "Calibri"
        sp4_5 = stf5.add_paragraph()
        sr4_5 = sp4_5.add_run(); sr4_5.text = "Evaluasi umur teknis unit & percepatan replacement unit yang tidak produktif."
        sr4_5.font.size = Pt(8.5); sr4_5.font.color.rgb = TEXT_DARK; sr4_5.font.name = "Calibri"


        # ================= SLIDE 6: KEY INSIGHTS \u2014 MAINTENANCE PER KATEGORI & DAMPAK BIAYA =================
        s = add_content_slide(f"KEY INSIGHTS \u2014 Maintenance per Kategori & Dampak Biaya s/d {period}", f"Analisis Frekuensi \u00b7 {snum6}{divisi_label}{kat_suffix}")

        # --- Data: Qty Pergantian & Total Biaya per kategori_sparepart ---
        qty_agg6 = pd.DataFrame(columns=["kategori_sparepart", "qty", "total_biaya"])
        if maint_data is not None and not maint_data.empty:
            m6 = maint_data.copy()
            if "lokasi" in m6.columns and site_list:
                m6 = m6[m6["lokasi"].isin(site_list)]
            if "bulan" in m6.columns and month_list:
                m6 = m6[m6["bulan"].isin(month_list)]
            kat_scope6 = set(data["kategori"].dropna().unique())
            if "kategori" in m6.columns:
                m6 = m6[m6["kategori"].isin(kat_scope6)]
            elif "nama_unit" in m6.columns:
                valid_units6 = set(data["nama_unit"].astype(str).str.strip().str.upper().unique())
                m6 = m6[m6["nama_unit"].astype(str).str.strip().str.upper().isin(valid_units6)]
            if "kategori_sparepart" in m6.columns:
                qty_agg6 = m6.groupby("kategori_sparepart", as_index=False).agg(
                    qty=("kategori_sparepart", "count"), total_biaya=("biaya", "sum"))
                qty_agg6 = qty_agg6.sort_values("total_biaya", ascending=False)

        panel_top6 = 1.0
        panel_bottom6 = 7.3
        panel_h6 = panel_bottom6 - panel_top6
        left_w6 = 6.4
        right_x6 = 7.0
        right_w6 = 5.9

        # ================= PANEL KIRI: TABEL Qty Pergantian per Kategori (dgn Status) =================
        add_card_panel(s, 0.4, panel_top6, left_w6, panel_h6, accent_color=NAVY)
        add_textbox(s, 0.55, panel_top6 + 0.08, left_w6 - 0.3, 0.3, f"Biaya Maintenance per Kategori \u2014 s/d {period}", size=13, bold=True, color=NAVY)

        tbl_top6 = panel_top6 + 0.45
        total_biaya_grand6 = qty_agg6["total_biaya"].sum() if not qty_agg6.empty else 0
        fair_share6 = (100 / len(qty_agg6)) if len(qty_agg6) else 0
        max_rows6 = 9
        qty_rows6 = []
        over_kategori6 = []
        for _, r in qty_agg6.head(max_rows6).iterrows():
            share_pct_row6 = (r["total_biaya"] / total_biaya_grand6 * 100) if total_biaya_grand6 else 0
            is_over = share_pct_row6 > fair_share6 * 1.3 if fair_share6 else False
            status_txt = f"\u2717 {share_pct_row6:.1f}%" if is_over else f"\u2713 {share_pct_row6:.1f}%"
            qty_rows6.append([str(r["kategori_sparepart"]), f"{int(r['qty']):,}", fmt_rp(r["total_biaya"]), status_txt])
            if is_over:
                over_kategori6.append(str(r["kategori_sparepart"]))
        if len(qty_agg6) > max_rows6:
            sisa6 = qty_agg6.iloc[max_rows6:]
            qty_rows6.append([f"+ {len(sisa6)} kategori lainnya", f"{int(sisa6['qty'].sum()):,}", fmt_rp(sisa6["total_biaya"].sum()), "-"])

        # Reservasi ruang tetap utk kotak catatan di bawah dulu, baru tabel mengisi sisa ruang yg ada
        note_h6 = 0.7
        tbl_h6_avail = panel_h6 - 0.45 - note_h6 - 0.25
        if qty_rows6:
            n_row6 = len(qty_rows6)
            row_h6 = max(0.22, min(0.4, tbl_h6_avail / (n_row6 + 1)))
            font6 = 9.5 if n_row6 <= 7 else (8.5 if n_row6 <= 9 else 7.5)
            add_table(s, 0.55, tbl_top6, left_w6 - 0.3, row_h6 * (n_row6 + 1),
                      ["Kategori Sparepart", "Qty", "Total Biaya", "Status"], qty_rows6,
                      col_widths=[2.5, 0.8, 1.5, 1.1], font_size=font6, header_size=font6)
        else:
            add_textbox(s, 0.55, tbl_top6 + 0.2, left_w6 - 0.3, 0.6, "Data Maintenance belum tersedia.", size=10, italic=True, color=TEXT_MUTED)

        # --- Catatan bawah panel kiri: kategori mana yang OVER rata-rata ---
        note_top6 = panel_bottom6 - note_h6 - 0.15
        if over_kategori6:
            over_txt6 = ", ".join(over_kategori6[:3]) + (", dll" if len(over_kategori6) > 3 else "")
            add_finding_box(s, 0.55, note_top6, left_w6 - 0.3, note_h6, "\u26A0",
                             f"Kategori dgn kontribusi biaya jauh di atas porsi wajar (rata-rata {fair_share6:.1f}% per kategori): {over_txt6}. "
                             f"Perlu investigasi penyebab tingginya "
                             f"frekuensi/biaya perbaikan pada kategori ini.",
                             RED_BG, RED, RED)
        elif qty_rows6:
            add_finding_box(s, 0.55, note_top6, left_w6 - 0.3, note_h6, "\u2705",
                             "Seluruh kategori sparepart berada dalam rentang biaya yang wajar (tidak ada outlier signifikan).",
                             GREEN_BG, GREEN, GREEN)

        # ================= PANEL KANAN ATAS: Analisis Kategori Tertinggi =================
        top_kat6 = qty_agg6.iloc[0] if not qty_agg6.empty else None
        box1_top6 = panel_top6
        box1_h6 = panel_h6 * 0.42

        box1_6 = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(right_x6), Inches(box1_top6), Inches(right_w6), Inches(box1_h6))
        box1_6.fill.solid(); box1_6.fill.fore_color.rgb = WHITE
        box1_6.line.color.rgb = BORDER; box1_6.line.width = Pt(0.75)
        box1_6.shadow.inherit = False
        strip1_6 = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(right_x6), Inches(box1_top6), Inches(right_w6), Inches(0.06))
        strip1_6.fill.solid(); strip1_6.fill.fore_color.rgb = RED
        strip1_6.line.fill.background(); strip1_6.shadow.inherit = False

        if top_kat6 is not None:
            biaya_per_kejadian6 = top_kat6["total_biaya"] / top_kat6["qty"] if top_kat6["qty"] else 0
            share_pct6 = (top_kat6["total_biaya"] / qty_agg6["total_biaya"].sum() * 100) if qty_agg6["total_biaya"].sum() else 0
            add_textbox(s, right_x6 + 0.15, box1_top6 + 0.14, right_w6 - 0.3, 0.3,
                        f"Analisis {top_kat6['kategori_sparepart']} ({top_kat6['qty']:.0f}x, {fmt_rp(top_kat6['total_biaya'])})",
                        size=11.5, bold=True, color=RED)
            bullets1_6 = [
                f"Kategori dengan KONTRIBUSI BIAYA TERTINGGI, menyumbang {share_pct6:.1f}% dari total biaya maintenance keseluruhan "
                f"(porsi wajar rata-rata: {fair_share6:.1f}% per kategori).",
                f"Rata-rata biaya per kejadian: {fmt_rp(biaya_per_kejadian6)} \u2014 {int(top_kat6['qty'])} kali kejadian s/d {period}.",
            ]
            if len(qty_agg6) > 1:
                kat2_6 = qty_agg6.iloc[1]
                selisih_kat6 = top_kat6["total_biaya"] - kat2_6["total_biaya"]
                bullets1_6.append(
                    f"Selisih dengan kategori tertinggi ke-2 ({kat2_6['kategori_sparepart']}, {fmt_rp(kat2_6['total_biaya'])}): "
                    f"{fmt_rp(selisih_kat6)} \u2014 menunjukkan kesenjangan yang perlu diperhatikan."
                )
            bullets1_6.append("Rekomendasi: telusuri riwayat kerusakan & evaluasi kesesuaian spare part dengan spesifikasi standar pabrikan.")
            bt6 = s.shapes.add_textbox(Inches(right_x6 + 0.15), Inches(box1_top6 + 0.5), Inches(right_w6 - 0.3), Inches(box1_h6 - 0.6))
            btf6 = bt6.text_frame; btf6.word_wrap = True
            for bi6, btxt6 in enumerate(bullets1_6):
                bpar6 = btf6.paragraphs[0] if bi6 == 0 else btf6.add_paragraph()
                bpar6.space_after = Pt(6)
                brun6 = bpar6.add_run(); brun6.text = f"●  {btxt6}"
                brun6.font.size = Pt(9.5); brun6.font.color.rgb = TEXT_DARK; brun6.font.name = "Calibri"
        else:
            add_textbox(s, right_x6 + 0.15, box1_top6 + 0.14, right_w6 - 0.3, box1_h6 - 0.3,
                        "Data Maintenance belum tersedia.", size=10, italic=True, color=TEXT_MUTED)

        # ================= PANEL KANAN BAWAH: Kalkulasi Dampak Biaya =================
        box2_top6 = box1_top6 + box1_h6 + 0.15
        box2_h6 = panel_h6 - box1_h6 - 0.15

        box2_6 = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(right_x6), Inches(box2_top6), Inches(right_w6), Inches(box2_h6))
        box2_6.fill.solid(); box2_6.fill.fore_color.rgb = WHITE
        box2_6.line.color.rgb = BORDER; box2_6.line.width = Pt(0.75)
        box2_6.shadow.inherit = False
        strip2_6 = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(right_x6), Inches(box2_top6), Inches(right_w6), Inches(0.06))
        strip2_6.fill.solid(); strip2_6.fill.fore_color.rgb = NAVY
        strip2_6.line.fill.background(); strip2_6.shadow.inherit = False
        add_textbox(s, right_x6 + 0.15, box2_top6 + 0.14, right_w6 - 0.3, 0.3, "Kalkulasi Dampak Biaya Maintenance", size=11.5, bold=True, color=NAVY)

        # Cari unit paling sering maintenance (utk kalkulasi dampak pendapatan)
        freq6 = pd.DataFrame()
        if maint_data is not None and not maint_data.empty and "nama_unit" in maint_data.columns:
            m6b = maint_data.copy()
            if "lokasi" in m6b.columns and site_list:
                m6b = m6b[m6b["lokasi"].isin(site_list)]
            if "bulan" in m6b.columns and month_list:
                m6b = m6b[m6b["bulan"].isin(month_list)]
            kat_scope6b = set(data["kategori"].dropna().unique())
            if "kategori" in m6b.columns:
                m6b = m6b[m6b["kategori"].isin(kat_scope6b)]
            elif "nama_unit" in m6b.columns:
                valid_units6b = set(data["nama_unit"].astype(str).str.strip().str.upper().unique())
                m6b = m6b[m6b["nama_unit"].astype(str).str.strip().str.upper().isin(valid_units6b)]
            freq6 = m6b.groupby(["lokasi", "nama_unit"], as_index=False).agg(
                jumlah_maintenance=("nama_unit", "count"), total_biaya=("biaya", "sum"))
            freq6 = freq6.sort_values("jumlah_maintenance", ascending=False)

        if not qty_agg6.empty:
            total_biaya_all6 = qty_agg6["total_biaya"].sum()
            bullets2_6 = [f"Total Biaya Maintenance (semua kategori): {fmt_rp(total_biaya_all6)} s/d {period}."]
            if top_kat6 is not None:
                bullets2_6.append(f"{top_kat6['kategori_sparepart']}: {fmt_rp(top_kat6['total_biaya'])} ({share_pct6:.1f}% dari total).")
            top3_kat6 = qty_agg6.head(3)
            top3_share6 = (top3_kat6["total_biaya"].sum() / total_biaya_all6 * 100) if total_biaya_all6 else 0
            top3_names6 = ", ".join(top3_kat6["kategori_sparepart"].tolist())
            bullets2_6.append(f"Top 3 kategori ({top3_names6}) menyumbang {top3_share6:.1f}% dari total biaya maintenance keseluruhan.")
            if not freq6.empty:
                top_unit6b = freq6.iloc[0]
                site_s6b = SITE_ABBR.get(top_unit6b["lokasi"], top_unit6b["lokasi"])
                unit_pend6b = data[data["nama_unit"].astype(str).str.strip().str.upper() == str(top_unit6b["nama_unit"]).strip().upper()]
                pend_r6b = unit_pend6b["pendapatan_realisasi"].sum() if not unit_pend6b.empty else None
                pend_b6b = unit_pend6b["pendapatan_budget"].sum() if not unit_pend6b.empty else None
                gap6b = (pend_r6b - pend_b6b) if (pend_r6b is not None and pend_b6b is not None) else None
                bullets2_6.append(f"Unit paling sering maintenance: {site_s6b} — {top_unit6b['nama_unit']} ({int(top_unit6b['jumlah_maintenance'])}x, {fmt_rp(top_unit6b['total_biaya'])}).")
                if gap6b is not None:
                    gap_sign6b = "MINUS" if gap6b < 0 else "PLUS"
                    bullets2_6.append(f"Dampak Pendapatan unit tsb: {gap_sign6b} {fmt_rp(abs(gap6b))} (Realisasi {fmt_rp(pend_r6b)} vs Budget {fmt_rp(pend_b6b)}).")

            bt2_6 = s.shapes.add_textbox(Inches(right_x6 + 0.15), Inches(box2_top6 + 0.5), Inches(right_w6 - 0.3), Inches(box2_h6 - 0.6))
            btf2_6 = bt2_6.text_frame; btf2_6.word_wrap = True
            for bi2_6, btxt2_6 in enumerate(bullets2_6):
                bpar2_6 = btf2_6.paragraphs[0] if bi2_6 == 0 else btf2_6.add_paragraph()
                bpar2_6.space_after = Pt(6)
                brun2_6 = bpar2_6.add_run(); brun2_6.text = f"●  {btxt2_6}"
                brun2_6.font.size = Pt(9); brun2_6.font.color.rgb = TEXT_DARK; brun2_6.font.name = "Calibri"
                # Highlight angka Rupiah dgn warna merah (segmen sesudah kata "Rp")
        else:
            add_textbox(s, right_x6 + 0.15, box2_top6 + 0.5, right_w6 - 0.3, box2_h6 - 0.6,
                        "Data Maintenance belum tersedia.", size=10, italic=True, color=TEXT_MUTED)




    # Jika filter kategori mencakup AB & TR sekaligus DAN keduanya benar-benar punya data, pisah jadi 2 blok:
    # TRANSPORTASI dulu (01-04), baru ALAT BERAT (05-08). Kalau salah satu kategori kosong (mis. Mining -> TR kosong
    # krn Tanjung digabung ke AB), render sebagai satu blok tunggal seperti biasa (tanpa penomoran 05-08 & label kosong).
    kat_set_render = set(kat_list) if kat_list else set()
    data_tr_check = data[data["kategori"] == "TR"]
    data_ab_check = data[data["kategori"] == "AB"]
    if {"AB", "TR"}.issubset(kat_set_render) and not data_tr_check.empty and not data_ab_check.empty:
        data_tr = data_tr_check.copy()
        data_ab = data_ab_check.copy()
        sm_tr = sasaran_mutu_data[sasaran_mutu_data["kategori"] == "TR"].copy() if (sasaran_mutu_data is not None and not sasaran_mutu_data.empty) else sasaran_mutu_data
        sm_ab = sasaran_mutu_data[sasaran_mutu_data["kategori"] == "AB"].copy() if (sasaran_mutu_data is not None and not sasaran_mutu_data.empty) else sasaran_mutu_data
        render_6_slides(data_tr, sm_tr, "01", "02", "03", "04", "05", "06", " · TRANSPORTASI")
        render_6_slides(data_ab, sm_ab, "07", "08", "09", "10", "11", "12", " · ALAT BERAT")
    else:
        render_6_slides(data, sasaran_mutu_data, "01", "02", "03", "04", "05", "06", "")

    buf = _io.BytesIO()
    prs.save(buf)
    return buf.getvalue()

colX, colY = st.columns([5, 1.4])
with colY:
    if st.button("📽️ Buat Presentasi (PPTX)", use_container_width=True, type="primary"):
        with st.spinner("Menyusun slide presentasi..."):
            maint_for_pptx = maint_df_site_bulan if not maint_raw.empty else pd.DataFrame()
            sparepart_for_pptx = sparepart_df_site_bulan if not sparepart_raw.empty else pd.DataFrame()
            pptx_bytes = build_pptx(df, maint_for_pptx, sparepart_for_pptx, sel_site, sel_month, sel_kat, sasaran_mutu_df)
        st.session_state["pptx_bytes"] = pptx_bytes
    if "pptx_bytes" in st.session_state:
        st.download_button(
            "⬇️ Unduh PPTX untuk RTM",
            data=st.session_state["pptx_bytes"],
            file_name="Laporan_Biaya_Pendapatan_BKMS.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            use_container_width=True,
        )

st.markdown("---")

# ---------------------------------------------------------------
# Klasifikasi satuan Rp/HM, Rp/KM, Rp/Tonase per baris data (dipakai di beberapa section)
# ---------------------------------------------------------------
def klasifikasi_satuan_prestasi(row):
    lok = row["lokasi"]
    kat = row["kategori"]
    if lok == "BUHUT LHL":
        return "Rp/Tonase"
    if lok in ("SUNGAI DANAU", "KUMAI"):
        if kat == "AB":
            return "Rp/HM"
        if kat == "TR":
            return "Rp/KM"
        return "Rp/HM"
    # TANJUNG, BUHUT, AMPAH, dan site lain di luar aturan eksplisit -> default Rp/HM
    return "Rp/HM"

df["satuan_prestasi"] = df.apply(klasifikasi_satuan_prestasi, axis=1)
satuan_suffix_map = {"Rp/HM": "/HM", "Rp/KM": "/KM", "Rp/Tonase": "/Ton"}

def rk_badge(pct, higher_is_better, na=False, small=False):
    cls_base = "rk-badge-sm" if small else "rk-badge"
    if na or pct is None:
        return f'<span class="{cls_base} rk-grey">N/A</span>'
    good = (pct >= 100) if higher_is_better else (pct <= 100)
    cls = "rk-green" if good else "rk-red"
    return f'<span class="{cls_base} {cls}">{pct:.1f}%</span>'

# ---------------------------------------------------------------
# 1. CAPAIAN PENDAPATAN  +  ANALISA PENYEBAB (di sampingnya)
# ---------------------------------------------------------------
if "bulan_no" in df.columns and not df.empty:
    bulan_terakhir_no = df["bulan_no"].max()
    bulan_terakhir_label = df.loc[df["bulan_no"] == bulan_terakhir_no, "bulan"].iloc[0]
    df_bulan_terakhir = df[df["bulan_no"] == bulan_terakhir_no]
else:
    bulan_terakhir_label = "-"
    df_bulan_terakhir = df

col_capaian, col_analisa = st.columns([1, 1.5])

with col_capaian:
    st.markdown('<h3 class="section-title">Capaian Pendapatan</h3>', unsafe_allow_html=True)

    pendapatan_pill, pendapatan_style = achievement_pill(ach_pendapatan, higher_is_better=True)

    st.markdown(kpi_card(
        icon="💰", icon_bg=RED, accent=RED,
        label="Pendapatan: Realisasi vs Target",
        value=fmt_rp(tot_pendapatan_r),
        budget_text=f"Target: {fmt_rp(tot_pendapatan_b)}",
        pill_text=pendapatan_pill, pill_style=pendapatan_style,
    ), unsafe_allow_html=True)

with col_analisa:
    st.markdown('<h3 class="section-title">Analisa: Penyebab Pendapatan Tidak Capai Budget</h3>', unsafe_allow_html=True)

    gap_rp = tot_pendapatan_b - tot_pendapatan_r
    rate_target_all = (tot_pendapatan_b / tot_prestasi_b) if tot_prestasi_b else None
    rate_realisasi_all = (tot_pendapatan_r / tot_prestasi_r) if tot_prestasi_r else None
    rate_ach = (rate_realisasi_all / rate_target_all * 100) if (rate_realisasi_all is not None and rate_target_all) else None

    if ach_pendapatan is not None and ach_pendapatan >= 100:
        st.markdown(
            f'<div class="insight-box">✅ Pendapatan sudah <b>mencapai/melampaui target</b> '
            f'({ach_pendapatan:.1f}%). Tidak ada gap yang perlu dianalisa lebih lanjut untuk periode/filter saat ini.</div>',
            unsafe_allow_html=True,
        )
    elif ach_pendapatan is None:
        st.markdown(
            '<div class="insight-box">⚠️ Target Pendapatan untuk kombinasi filter ini adalah <b>Rp 0</b>, '
            'sehingga capaian tidak dapat dihitung. Analisa penyebab gap tidak berlaku untuk filter saat ini.</div>',
            unsafe_allow_html=True,
        )
    else:
        drivers = []
        if pct_populasi is not None:
            drivers.append(("Populasi (jumlah unit beroperasi)", pct_populasi))
        if ach_prestasi is not None:
            drivers.append(("Prestasi (volume pekerjaan)", ach_prestasi))
        if rate_ach is not None:
            drivers.append(("Tarif/Rate Rp per satuan", rate_ach))

        penyebab_utama = min(drivers, key=lambda x: x[1]) if drivers else None

        poin = []
        poin.append(f"Pendapatan hanya mencapai <b>{ach_pendapatan:.1f}%</b> dari target (gap <b>{fmt_rp(gap_rp)}</b>).")

        if target_populasi:
            gap_pop = target_populasi - realisasi_populasi
            if gap_pop > 0:
                poin.append(
                    f"<b>Populasi:</b> dari <b>{target_populasi}</b> unit target pada bulan {bulan_terakhir_label}, "
                    f"hanya <b>{realisasi_populasi}</b> unit yang mencatatkan realisasi pendapatan — "
                    f"<b>{gap_pop} unit ({100 - pct_populasi:.1f}%) tidak beroperasi/tidak menghasilkan pendapatan sama sekali</b>."
                )
            else:
                poin.append(f"<b>Populasi:</b> seluruh unit target ({target_populasi} unit) sudah beroperasi dan mencatatkan pendapatan.")

        if ach_prestasi is not None:
            if ach_prestasi < 100:
                poin.append(f"<b>Prestasi</b> (volume pekerjaan unit yang beroperasi) juga di bawah target, hanya <b>{ach_prestasi:.1f}%</b>.")
            else:
                poin.append(f"<b>Prestasi</b> dari unit yang beroperasi sebenarnya sudah tercapai (<b>{ach_prestasi:.1f}%</b>) — bukan penyebab utama gap.")

        if rate_ach is not None:
            if rate_ach < 100:
                poin.append(f"<b>Tarif/Rate</b> Rp per satuan prestasi realisasi ({fmt_rp(rate_realisasi_all)}) juga lebih rendah dari target ({fmt_rp(rate_target_all)}), yaitu <b>{rate_ach:.1f}%</b>.")
            else:
                poin.append(f"<b>Tarif/Rate</b> Rp per satuan prestasi sudah sesuai/di atas target (<b>{rate_ach:.1f}%</b>) — bukan penyebab utama gap.")

        # Jenis Unit mana yang paling bermasalah (kontribusi gap pendapatan terbesar)
        ju_problem = df.groupby("jenis_unit", as_index=False).agg(
            target_pendapatan=("pendapatan_budget", "sum"),
            realisasi_pendapatan=("pendapatan_realisasi", "sum"),
            target_prestasi=("prestasi_budget", "sum"),
            realisasi_prestasi=("prestasi_realisasi", "sum"),
        )
        ju_problem["gap"] = ju_problem["target_pendapatan"] - ju_problem["realisasi_pendapatan"]
        ju_problem["capaian"] = ju_problem.apply(
            lambda r: (r["realisasi_pendapatan"] / r["target_pendapatan"] * 100) if r["target_pendapatan"] else None, axis=1
        )
        ju_problem["capaian_prestasi"] = ju_problem.apply(
            lambda r: (r["realisasi_prestasi"] / r["target_prestasi"] * 100) if r["target_prestasi"] else None, axis=1
        )
        ju_bermasalah = ju_problem[(ju_problem["gap"] > 0) & (ju_problem["capaian"].notna()) & (ju_problem["capaian"] < 100)]
        ju_bermasalah = ju_bermasalah.sort_values("gap", ascending=False).head(5)

        if not ju_bermasalah.empty:
            list_items = ""
            for _, jr in ju_bermasalah.iterrows():
                cap_prestasi_txt = f", prestasi {jr['capaian_prestasi']:.1f}%" if pd.notna(jr["capaian_prestasi"]) else ""
                list_items += (
                    f"<li style='margin-bottom:3px;'><b>{jr['jenis_unit']}</b> — capaian pendapatan {jr['capaian']:.1f}%"
                    f"{cap_prestasi_txt}, gap {fmt_rp(jr['gap'])}</li>"
                )
            poin.append(
                f"<b>🚩 Jenis Unit paling bermasalah</b> (kontribusi gap pendapatan terbesar):"
                f"<ul style='padding-left:16px; margin:4px 0 0 0;'>{list_items}</ul>"
            )

        if penyebab_utama:
            poin.append(f"➡️ <b>Penyebab paling dominan:</b> {penyebab_utama[0]} (capaian terendah, {penyebab_utama[1]:.1f}%).")

        bullets_html = "".join([f"<li style='margin-bottom:6px;'>{p}</li>" for p in poin])
        st.markdown(f'<div class="insight-box"><ul style="padding-left:18px; margin:0;">{bullets_html}</ul></div>', unsafe_allow_html=True)

st.markdown("---")

# ---------------------------------------------------------------
# RINGKASAN PER JENIS UNIT
# ---------------------------------------------------------------
st.markdown('<h3 class="section-title">Ringkasan per Jenis Unit</h3>', unsafe_allow_html=True)

ju_target_pop = df_bulan_terakhir[df_bulan_terakhir["pendapatan_budget"] > 0].groupby("jenis_unit")["nama_unit"].nunique()
ju_real_pop = df_bulan_terakhir[df_bulan_terakhir["pendapatan_realisasi"] > 0].groupby("jenis_unit")["nama_unit"].nunique()
ju_sat_mode = df.groupby("jenis_unit")["satuan_prestasi"].agg(lambda s: s.mode().iat[0] if not s.mode().empty else "Rp/HM")

ju = df.groupby("jenis_unit", as_index=False).agg(
    target_pendapatan=("pendapatan_budget", "sum"),
    realisasi_pendapatan=("pendapatan_realisasi", "sum"),
    target_prestasi=("prestasi_budget", "sum"),
    realisasi_prestasi=("prestasi_realisasi", "sum"),
)
ju["target_populasi"] = ju["jenis_unit"].map(ju_target_pop).fillna(0).astype(int)
ju["realisasi_populasi"] = ju["jenis_unit"].map(ju_real_pop).fillna(0).astype(int)
ju["satuan"] = ju["jenis_unit"].map(ju_sat_mode).fillna("Rp/HM")
ju["suffix"] = ju["satuan"].map(satuan_suffix_map)
ju["rate_target"] = ju.apply(lambda r: (r["target_pendapatan"] / r["target_prestasi"]) if r["target_prestasi"] else None, axis=1)
ju["rate_realisasi"] = ju.apply(lambda r: (r["realisasi_pendapatan"] / r["realisasi_prestasi"]) if r["realisasi_prestasi"] else None, axis=1)
ju["capaian"] = ju.apply(lambda r: (r["realisasi_pendapatan"] / r["target_pendapatan"] * 100) if r["target_pendapatan"] else None, axis=1)
ju["selisih_pendapatan"] = ju["realisasi_pendapatan"] - ju["target_pendapatan"]
ju = ju.sort_values("selisih_pendapatan", ascending=False)

def fmt_rp_signed(x, pct=None):
    if x is None or pd.isna(x):
        return "-"
    sign = "▲ +" if x > 0 else ("▼ -" if x < 0 else "")
    base = f"{sign}{fmt_rp(abs(x))}" if x != 0 else fmt_rp(0)
    if pct is not None and pd.notna(pct):
        pct_sign = "+" if pct > 0 else ("" if pct < 0 else "+")
        base += f" ({pct_sign}{pct:.1f}%)"
    return base

ju["selisih_pct"] = ju.apply(
    lambda r: (r["selisih_pendapatan"] / r["target_pendapatan"] * 100) if r["target_pendapatan"] else None, axis=1
)

show_ju = pd.DataFrame({
    "Jenis Unit": ju["jenis_unit"],
    "Target Populasi": ju["target_populasi"],
    "Realisasi Populasi": ju["realisasi_populasi"],
    "Target Pendapatan": ju["target_pendapatan"].apply(fmt_rp),
    "Realisasi Pendapatan": ju["realisasi_pendapatan"].apply(fmt_rp),
    "Selisih (Realisasi − Target)": ju.apply(lambda r: fmt_rp_signed(r["selisih_pendapatan"], r["selisih_pct"]), axis=1),
    "Target Prestasi": ju["target_prestasi"].apply(lambda v: f"{v:,.0f}"),
    "Realisasi Prestasi": ju["realisasi_prestasi"].apply(lambda v: f"{v:,.0f}"),
    "Target (Rp/Satuan)": ju.apply(lambda r: f"{fmt_rp(r['rate_target'])}{r['suffix']}" if pd.notna(r["rate_target"]) else "-", axis=1),
    "Realisasi (Rp/Satuan)": ju.apply(lambda r: f"{fmt_rp(r['rate_realisasi'])}{r['suffix']}" if pd.notna(r["rate_realisasi"]) else "-", axis=1),
    "Capaian (%)": ju["capaian"],
})

def _color_selisih(val):
    if isinstance(val, str):
        if val.startswith("▲"):
            return f"color: {CHART_GREEN}; font-weight: 700;"
        if val.startswith("▼"):
            return f"color: {RED}; font-weight: 700;"
    return ""

try:
    styled_ju = show_ju.style.map(_color_selisih, subset=["Selisih (Realisasi − Target)"])
except AttributeError:
    styled_ju = show_ju.style.applymap(_color_selisih, subset=["Selisih (Realisasi − Target)"])

st.dataframe(
    styled_ju,
    use_container_width=True,
    hide_index=True,
    height=480,
    column_config={
        "Capaian (%)": st.column_config.ProgressColumn(
            "Capaian (%)", format="%.1f%%", min_value=0, max_value=150,
        ),
    },
)

st.markdown("---")

# ---------------------------------------------------------------
# CAPAIAN PRESTASI PER SATUAN (Rp/HM, Rp/KM, Rp/Tonase)
# ---------------------------------------------------------------
st.markdown('<h3 class="section-title">Capaian Prestasi per Satuan</h3>', unsafe_allow_html=True)

satuan_order = ["Rp/HM", "Rp/KM", "Rp/Tonase"]
satuan_icon = {"Rp/HM": ("⏱️", CHART_GREEN), "Rp/KM": ("🚚", GOLD), "Rp/Tonase": ("⚖️", RED)}
satuan_suffix = {"Rp/HM": "/HM", "Rp/KM": "/KM", "Rp/Tonase": "/Ton"}
cols_satuan = st.columns(3)
for i, sat in enumerate(satuan_order):
    sub = df[df["satuan_prestasi"] == sat]
    pendapatan_r = sub["pendapatan_realisasi"].sum()
    pendapatan_b = sub["pendapatan_budget"].sum()
    prestasi_r = sub["prestasi_realisasi"].sum()
    prestasi_b = sub["prestasi_budget"].sum()
    rate_r = (pendapatan_r / prestasi_r) if prestasi_r else None
    rate_b = (pendapatan_b / prestasi_b) if prestasi_b else None
    ach = (rate_r / rate_b * 100) if (rate_r is not None and rate_b) else None
    pill_txt, pill_style = achievement_pill(ach, higher_is_better=True)
    icon, icon_color = satuan_icon[sat]
    suf = satuan_suffix[sat]
    with cols_satuan[i]:
        st.markdown(kpi_card(
            icon=icon, icon_bg=icon_color, accent=icon_color,
            label=f"Capaian {sat} (Pendapatan ÷ Prestasi)",
            value=(f"{fmt_rp(rate_r)}{suf}" if rate_r is not None else "-"),
            budget_text=(f"Target: {fmt_rp(rate_b)}{suf} • {sub['nama_unit'].nunique():,} unit" if rate_b is not None else f"{sub['nama_unit'].nunique():,} unit"),
            pill_text=pill_txt, pill_style=pill_style,
        ), unsafe_allow_html=True)

st.markdown("---")

# ---------------------------------------------------------------
# CAPAIAN BIAYA PER SATUAN (Rp/HM, Rp/KM, Rp/Tonase) — pembagian sama dengan Prestasi
# ---------------------------------------------------------------
st.markdown('<h3 class="section-title">Capaian Biaya per Satuan</h3>', unsafe_allow_html=True)

cols_satuan_biaya = st.columns(3)
for i, sat in enumerate(satuan_order):
    sub = df[df["satuan_prestasi"] == sat]
    biaya_r = sub["total_biaya_realisasi"].sum()
    biaya_b = sub["total_biaya_budget"].sum()
    prestasi_r = sub["prestasi_realisasi"].sum()
    prestasi_b = sub["prestasi_budget"].sum()
    rate_r = (biaya_r / prestasi_r) if prestasi_r else None
    rate_b = (biaya_b / prestasi_b) if prestasi_b else None
    ach = (rate_r / rate_b * 100) if (rate_r is not None and rate_b) else None
    pill_txt, pill_style = achievement_pill(ach, higher_is_better=False)
    icon, icon_color = satuan_icon[sat]
    suf = satuan_suffix[sat]
    with cols_satuan_biaya[i]:
        st.markdown(kpi_card(
            icon=icon, icon_bg=icon_color, accent=icon_color,
            label=f"Capaian Biaya {sat} (Biaya ÷ Prestasi)",
            value=(f"{fmt_rp(rate_r)}{suf}" if rate_r is not None else "-"),
            budget_text=(f"Target: {fmt_rp(rate_b)}{suf} • {sub['nama_unit'].nunique():,} unit" if rate_b is not None else f"{sub['nama_unit'].nunique():,} unit"),
            pill_text=pill_txt, pill_style=pill_style,
        ), unsafe_allow_html=True)

st.markdown("---")

# ---------------------------------------------------------------
# RINGKASAN BIAYA (tabel Budget vs Aktual vs Capaian)
# ---------------------------------------------------------------
st.markdown('<h3 class="section-title">Ringkasan Biaya</h3>', unsafe_allow_html=True)

tot_prestasi_r_all = df["prestasi_realisasi"].sum()
tot_prestasi_b_all = df["prestasi_budget"].sum()
ach_prestasi_all = (tot_prestasi_r_all / tot_prestasi_b_all * 100) if tot_prestasi_b_all else None
tot_qty_bbm_b_all = df["qty_bbm_budget"].sum()
tot_qty_bbm_r_all = df["qty_bbm_realisasi"].sum()

if sel_kat == ["TR"]:
    kat_label_suffix = "/KM"
elif sel_kat == ["AB"]:
    kat_label_suffix = "/satuan"
else:
    kat_label_suffix = ""

def biaya_row(label, real_col, budget_col, is_bbm=False, raw=False):
    comp_r_raw = df[real_col].sum()
    comp_b_raw = df[budget_col].sum()

    if raw:
        rate_b = comp_b_raw
        rate_r = comp_r_raw
        suffix = ""
        display_label = label
    elif is_bbm:
        rate_b = (comp_b_raw / tot_qty_bbm_b_all) if tot_qty_bbm_b_all else None
        rate_r = (comp_r_raw / tot_qty_bbm_r_all) if tot_qty_bbm_r_all else None
        suffix = "/Ltr"
        display_label = f"{label}/Ltr"
    else:
        rate_b = (comp_b_raw / tot_prestasi_b_all) if tot_prestasi_b_all else None
        rate_r = (comp_r_raw / tot_prestasi_r_all) if tot_prestasi_r_all else None
        suffix = ""
        display_label = f"{label}{kat_label_suffix}"

    ach = (rate_r / rate_b * 100) if (rate_r is not None and rate_b) else None
    return dict(label=display_label, budget=rate_b, aktual=rate_r, suffix=suffix, ach=ach,
                capaian_prestasi=ach_prestasi_all, aktual_raw=comp_r_raw)

ringkasan_rows = [
    biaya_row("Total Biaya", "total_biaya_realisasi", "total_biaya_budget", raw=True),
    biaya_row("Upah Operator", "upah_realisasi", "upah_budget"),
    biaya_row("Biaya BBM", "biaya_bbm_realisasi", "biaya_bbm_budget", is_bbm=True),
    biaya_row("Biaya Maintenance", "maintenance_realisasi", "maintenance_budget"),
    biaya_row("Penyusutan", "penyusutan_realisasi", "penyusutan_budget"),
    biaya_row("Lainnya", "lainnya_realisasi", "lainnya_budget"),
    biaya_row("Biaya Tidak Langsung", "biaya_tidak_langsung_realisasi", "biaya_tidak_langsung_budget"),
]

table_rows_html = ""
for row in ringkasan_rows:
    budget_disp = f"{fmt_rp(row['budget'])}{row['suffix']}" if row["budget"] is not None else "-"
    aktual_disp = f"{fmt_rp(row['aktual'])}{row['suffix']}" if row["aktual"] is not None else "-"
    hide_cp = row["label"] in ("Total Biaya", "Biaya Tidak Langsung")
    table_rows_html += f"""
    <tr>
        <td>{row['label']}</td>
        <td>{budget_disp}</td>
        <td>{aktual_disp}</td>
        <td>{rk_badge(row['ach'], higher_is_better=False, na=(row['ach'] is None), small=True)}</td>
        <td>{rk_badge(row['capaian_prestasi'], higher_is_better=True, na=(hide_cp or row['capaian_prestasi'] is None), small=True)}</td>
    </tr>"""

col_tbl, col_pie = st.columns([1, 1])
with col_tbl:
    st.markdown("##### Ringkasan Biaya (Budget vs Aktual)")
    st.markdown(f"""
    <table class="ringkasan-table-sm">
        <thead>
            <tr><th>Metrik</th><th>Budget</th><th>Aktual</th><th>Capaian</th><th>Capaian Prestasi</th></tr>
        </thead>
        <tbody>{table_rows_html}
        </tbody>
    </table>
    """, unsafe_allow_html=True)
with col_pie:
    st.markdown("##### Komposisi Biaya Aktual terhadap Total Biaya")
    comp_pie_df = pd.DataFrame({
        "Komponen": [r["label"] for r in ringkasan_rows if r["label"] != "Total Biaya"],
        "Biaya": [r["aktual_raw"] for r in ringkasan_rows if r["label"] != "Total Biaya"],
    })
    comp_pie_df = comp_pie_df[comp_pie_df["Biaya"] > 0]
    if not comp_pie_df.empty:
        fig_comp = px.pie(comp_pie_df, names="Komponen", values="Biaya", hole=0.5)
        fig_comp.update_traces(textinfo="percent", textfont=dict(size=11))
        fig_comp.update_layout(
            height=360, margin=dict(t=8, b=8, l=8, r=8),
            showlegend=True, legend=dict(font=dict(size=10)),
        )
        st.plotly_chart(style_fig(fig_comp), use_container_width=True)
    else:
        st.info("Data biaya aktual belum tersedia untuk ditampilkan sebagai diagram komposisi.")

st.markdown("---")

# ---------------------------------------------------------------
# 4. REKAP BIAYA MAINTENANCE (filter ID Unit + Rutin/Non Rutin + frekuensi)
# ---------------------------------------------------------------
st.markdown('<h3 class="section-title">Rekap Biaya Maintenance</h3>', unsafe_allow_html=True)

sel_unit_maint = []

if maint_raw.empty:
    st.info("Data maintenance belum tersedia. Upload file Pemeliharaan (.xls/.xlsx) di sidebar untuk menampilkan bagian ini.")
else:
    maint_df_site_bulan = maint_df_site_bulan.copy()
    maint_df_site_bulan["unit_label"] = maint_df_site_bulan["nama_unit"].apply(_unit_label)
    unit_maint_opts = sorted(maint_df_site_bulan["unit_label"].dropna().unique().tolist())

    # Rekonsiliasi: Total Maintenance = Pemakaian Persediaan (Sparepart) + Service Luar (residual)
    total_maint_all = maint_df_site_bulan["biaya"].sum()
    total_persediaan_all = sparepart_df_site_bulan["biaya"].sum() if not sparepart_raw.empty else 0
    service_luar_all = total_maint_all - total_persediaan_all
    pct_persediaan = (total_persediaan_all / total_maint_all * 100) if total_maint_all else 0
    pct_service_luar = (service_luar_all / total_maint_all * 100) if total_maint_all else 0

    st.markdown("##### Rekonsiliasi: Total Maintenance = Pemakaian Persediaan + Service Luar")
    rc1, rc2, rc3 = st.columns(3)
    with rc1:
        st.markdown(kpi_card(
            icon="🔧", icon_bg=GREY, accent=GREY,
            label="Total Biaya Maintenance", value=fmt_rp(total_maint_all),
            budget_text="Persediaan + Service Luar",
            pill_text="100%", pill_style="kpi-pill-amber",
        ), unsafe_allow_html=True)
    with rc2:
        st.markdown(kpi_card(
            icon="📦", icon_bg=CHART_GREEN, accent=CHART_GREEN,
            label="Pemakaian Persediaan (Sparepart)", value=fmt_rp(total_persediaan_all),
            budget_text="Dari data Rincian Pemakaian",
            pill_text=f"{pct_persediaan:.1f}% dari total", pill_style="kpi-pill-green",
        ), unsafe_allow_html=True)
    with rc3:
        st.markdown(kpi_card(
            icon="🛠️", icon_bg=GOLD, accent=GOLD,
            label="Service Luar (di luar persediaan)", value=fmt_rp(service_luar_all),
            budget_text="Selisih Total − Persediaan",
            pill_text=f"{pct_service_luar:.1f}% dari total", pill_style="kpi-pill-amber",
        ), unsafe_allow_html=True)

    st.markdown("---")

    sel_unit_maint = st.multiselect(
        "Filter berdasarkan ID Unit (opsional, kosongkan = semua unit) — ketik ID Unit atau nama unit",
        unit_maint_opts, default=[],
    )

    maint_df = maint_df_site_bulan
    if sel_unit_maint:
        maint_df = maint_df[maint_df["unit_label"].isin(sel_unit_maint)]

    if maint_df.empty:
        st.warning("Tidak ada data maintenance untuk kombinasi filter yang dipilih.")
    else:
        n_transaksi = len(maint_df)
        rutin_biaya = maint_df.loc[maint_df["jenis_pemeliharaan"] == "RUTIN", "biaya"].sum()
        nonrutin_biaya = maint_df.loc[maint_df["jenis_pemeliharaan"] == "NON RUTIN", "biaya"].sum()
        rutin_n = int((maint_df["jenis_pemeliharaan"] == "RUTIN").sum())
        nonrutin_n = int((maint_df["jenis_pemeliharaan"] == "NON RUTIN").sum())

        # Total Biaya Maintenance = total biaya di data_maintenance.csv (maint_df, sudah difilter ID Unit)
        total_maint = maint_df["biaya"].sum()

        # Workshop Sendiri = total biaya di data_sparepart.csv (pemakaian persediaan), scope ID Unit yang sama
        if not sparepart_raw.empty:
            sparepart_scope = sparepart_df_site_bulan.copy()
            sparepart_scope["unit_label"] = sparepart_scope["nama_unit"].apply(_unit_label)
            if sel_unit_maint:
                sparepart_scope = sparepart_scope[sparepart_scope["unit_label"].isin(sel_unit_maint)]
        else:
            sparepart_scope = pd.DataFrame(columns=["biaya"])
        workshop_biaya = sparepart_scope["biaya"].sum() if not sparepart_scope.empty else 0
        workshop_n = len(sparepart_scope)

        # Service Luar = selisih Total Biaya Maintenance (data_maintenance.csv) - Workshop Sendiri (data_sparepart.csv)
        service_luar_biaya = total_maint - workshop_biaya

        # Budget Maintenance diambil dari kolom maintenance_budget di data_bkms.csv, dengan
        # scope site/bulan/kategori yang sama, dan ikut mengikuti filter ID Unit jika dipilih.
        if sel_unit_maint:
            budget_scope = df[df["nama_unit"].apply(_unit_label).isin(sel_unit_maint)]
        else:
            budget_scope = df
        total_maint_budget = budget_scope["maintenance_budget"].sum()
        maint_ach = achievement(total_maint, total_maint_budget)
        if maint_ach is not None:
            diff_pp = maint_ach - 100  # positif = over budget (buruk), negatif = under budget (baik)
            maint_delta = f"{diff_pp:+.1f}% vs Budget {fmt_rp(total_maint_budget)}"
        else:
            maint_delta = "Budget belum tersedia"

        k1, k2, k3, k4, k5 = st.columns(5)
        k1.metric("Total Biaya Maintenance", fmt_rp(total_maint), maint_delta, delta_color="inverse")
        k2.metric("Workshop Sendiri", fmt_rp(workshop_biaya))
        k3.metric("Service Luar", fmt_rp(service_luar_biaya))
        k4.metric("Biaya Rutin", fmt_rp(rutin_biaya))
        k5.metric("Biaya Non Rutin", fmt_rp(nonrutin_biaya))

        rekap = maint_df.groupby(["kategori_sparepart", "jenis_pemeliharaan"], as_index=False).agg(
            jumlah_transaksi=("biaya", "count"),
            total_biaya=("biaya", "sum"),
        )
        rekap["total_biaya_jt"] = rekap["total_biaya"] / 1e6
        rekap["label_biaya"] = rekap["total_biaya"].apply(fmt_rp)

        # Urutkan kategori berdasarkan total biaya (terbesar di atas) supaya chart menyamping rapi
        kategori_order = (
            rekap.groupby("kategori_sparepart")["total_biaya"].sum().sort_values(ascending=True).index.tolist()
        )

        fig_rekap = px.bar(
            rekap, y="kategori_sparepart", x="total_biaya_jt", color="jenis_pemeliharaan",
            orientation="h", barmode="group", color_discrete_map={"RUTIN": CHART_GREEN, "NON RUTIN": GOLD},
            labels={"kategori_sparepart": "Kategori Sparepart / Sistem", "total_biaya_jt": "Total Biaya (Juta Rupiah)", "jenis_pemeliharaan": "Jenis"},
            text="label_biaya",
            category_orders={"kategori_sparepart": kategori_order},
        )
        fig_rekap.update_traces(textposition="outside", textfont=dict(size=10, color=TEXT_LIGHT))
        fig_rekap.update_layout(title="Maintenance atas Apa Saja — Rutin vs Non Rutin", height=550,
                                 legend=dict(orientation="h", y=1.08), margin=dict(t=60, b=10, l=10))
        fig_rekap.update_xaxes(ticksuffix=" Jt")
        st.plotly_chart(style_fig(fig_rekap), use_container_width=True)

        st.markdown("##### Rincian: Kategori, Jenis, Frekuensi (Berapa Kali), Total Biaya")
        rekap_tbl = rekap.sort_values("total_biaya", ascending=False)[
            ["kategori_sparepart", "jenis_pemeliharaan", "jumlah_transaksi", "total_biaya"]
        ].rename(columns={
            "kategori_sparepart": "Kategori (Maintenance atas Apa Saja)",
            "jenis_pemeliharaan": "Jenis (Rutin / Non Rutin)",
            "jumlah_transaksi": "Berapa Kali (Jumlah Transaksi)",
            "total_biaya": "Total Biaya",
        })
        total_row = pd.DataFrame([{
            "Kategori (Maintenance atas Apa Saja)": "TOTAL",
            "Jenis (Rutin / Non Rutin)": "",
            "Berapa Kali (Jumlah Transaksi)": rekap_tbl["Berapa Kali (Jumlah Transaksi)"].sum(),
            "Total Biaya": rekap_tbl["Total Biaya"].sum(),
        }])
        rekap_tbl_display = pd.concat([rekap_tbl, total_row], ignore_index=True)
        st.dataframe(
            rekap_tbl_display, use_container_width=True, hide_index=True, height=420,
            column_config={"Total Biaya": st.column_config.NumberColumn(format="Rp %,.0f")},
        )
        csv_rekap = rekap_tbl_display.to_csv(index=False).encode("utf-8")
        st.download_button("⬇️ Unduh Rekap Maintenance (CSV)", csv_rekap, file_name="rekap_maintenance_bkms.csv", mime="text/csv")

st.markdown("---")

# ---------------------------------------------------------------
# REKAP PEMAKAIAN SPAREPART (PERSEDIAAN)
# ---------------------------------------------------------------
st.markdown('<h3 class="section-title">Rekap Pemakaian Sparepart (Persediaan)</h3>', unsafe_allow_html=True)

if sparepart_raw.empty:
    st.info("Data pemakaian sparepart belum tersedia. Upload file Rincian Pemakaian (.xls/.xlsx) di sidebar untuk menampilkan bagian ini.")
else:
    sparepart_df_site_bulan = sparepart_df_site_bulan.copy()
    sparepart_df_site_bulan["unit_label"] = sparepart_df_site_bulan["nama_unit"].apply(_unit_label)

    sparepart_df = sparepart_df_site_bulan
    if sel_unit_maint:
        sparepart_df = sparepart_df[sparepart_df["unit_label"].isin(sel_unit_maint)]

    if sel_unit_maint:
        st.caption(f"🔗 Sedang difilter mengikuti ID Unit: {', '.join(sel_unit_maint)}")

    if sparepart_df.empty:
        st.warning("Tidak ada data pemakaian sparepart untuk kombinasi filter yang dipilih.")
    else:
        total_sp = sparepart_df["biaya"].sum()
        total_qty_trx = len(sparepart_df)
        n_jenis_barang = sparepart_df["nama_barang"].nunique()

        s1, s2, s3 = st.columns(3)
        s1.metric("Total Biaya Pemakaian Sparepart", fmt_rp(total_sp))
        s2.metric("Jumlah Transaksi Pengambilan", f"{total_qty_trx:,}")
        s3.metric("Jenis Barang Berbeda", f"{n_jenis_barang:,}")

        colS1, colS2 = st.columns([3, 2])

        # Peta warna konsisten per kategori sparepart, dipakai bersama oleh chart Top 15 Barang & pie chart
        all_cats_sp = sorted(sparepart_df["kategori_sparepart"].dropna().unique().tolist())
        _palette_sp = px.colors.qualitative.Plotly
        color_map_sp = {cat: _palette_sp[i % len(_palette_sp)] for i, cat in enumerate(all_cats_sp)}

        with colS1:
            top_barang = sparepart_df.groupby("nama_barang", as_index=False).agg(
                total_qty=("qty", "sum"), total_biaya=("biaya", "sum"),
            ).sort_values("total_biaya", ascending=False).head(15)

            # Tentukan kategori sparepart dominan (berdasarkan biaya terbesar) untuk tiap barang
            item_cat = sparepart_df.groupby(["nama_barang", "kategori_sparepart"], as_index=False)["biaya"].sum()
            item_cat = item_cat.sort_values("biaya", ascending=False).drop_duplicates("nama_barang")
            top_barang = top_barang.merge(item_cat[["nama_barang", "kategori_sparepart"]], on="nama_barang", how="left")

            top_barang["total_biaya_jt"] = top_barang["total_biaya"] / 1e6
            top_barang["label_biaya"] = top_barang["total_biaya"].apply(fmt_rp)

            fig_sp1 = px.bar(
                top_barang, y="nama_barang", x="total_biaya_jt", color="kategori_sparepart",
                orientation="h", color_discrete_map=color_map_sp, text="label_biaya",
                labels={"nama_barang": "Barang", "total_biaya_jt": "Juta Rupiah", "kategori_sparepart": "Kategori"},
            )
            fig_sp1.update_traces(textposition="outside", textfont=dict(size=10, color=TEXT_LIGHT))
            fig_sp1.update_layout(title="Top 15 Barang berdasarkan Biaya", xaxis_title="Juta Rupiah",
                                   height=560, margin=dict(t=60, b=10, l=10, r=70),
                                   legend=dict(font=dict(size=9), title="Kategori"))
            fig_sp1.update_xaxes(ticksuffix=" Jt")
            fig_sp1.update_yaxes(autorange="reversed")
            st.plotly_chart(style_fig(fig_sp1), use_container_width=True)

        with colS2:
            cat_sp = sparepart_df.groupby("kategori_sparepart", as_index=False)["biaya"].sum()
            fig_sp2 = px.pie(cat_sp, names="kategori_sparepart", values="biaya", hole=0.5,
                              color="kategori_sparepart", color_discrete_map=color_map_sp)
            fig_sp2.update_layout(title="Komposisi per Kategori Sparepart", height=460, margin=dict(t=60, b=10),
                                   showlegend=True, legend=dict(font=dict(size=9)))
            st.plotly_chart(style_fig(fig_sp2), use_container_width=True)

        st.markdown("##### Rincian Pemakaian per Barang")
        search_sp = st.text_input("🔍 Cari nama barang / part number...", "", key="search_sparepart")
        show_sp = sparepart_df[["tanggal", "lokasi", "nama_unit", "kategori_sparepart", "jenis_pemeliharaan",
                                 "kode_barang", "part_number", "nama_barang", "qty", "satuan", "biaya"]].rename(columns={
            "tanggal": "Tanggal", "lokasi": "Site", "nama_unit": "Nama Unit", "kategori_sparepart": "Kategori Sparepart",
            "jenis_pemeliharaan": "Jenis", "kode_barang": "Kode Barang", "part_number": "Part Number",
            "nama_barang": "Nama Barang", "qty": "Qty", "satuan": "Satuan", "biaya": "Biaya",
        })
        if search_sp:
            mask = (show_sp["Nama Barang"].str.contains(search_sp, case=False, na=False) |
                    show_sp["Part Number"].str.contains(search_sp, case=False, na=False))
            show_sp = show_sp[mask]
        st.dataframe(
            show_sp.sort_values("Biaya", ascending=False),
            use_container_width=True, height=380,
            column_config={"Biaya": st.column_config.NumberColumn(format="Rp %,.0f")},
        )
        csv_sp = show_sp.to_csv(index=False).encode("utf-8")
        st.download_button("⬇️ Unduh Rincian Pemakaian Sparepart (CSV)", csv_sp, file_name="rincian_sparepart_bkms.csv", mime="text/csv")

st.markdown("---")

# ---------------------------------------------------------------
# 6. ANALISA: PENYEBAB CAPAIAN PENDAPATAN
# ---------------------------------------------------------------
st.markdown('<h3 class="section-title">Analisa: Penyebab Capaian Pendapatan</h3>', unsafe_allow_html=True)

if ach_pendapatan is None:
    st.info("Target Pendapatan belum tersedia untuk kombinasi filter ini, sehingga analisa capaian tidak bisa dihitung.")
elif ach_pendapatan >= 100:
    st.success(f"Realisasi Pendapatan sudah mencapai **{ach_pendapatan:.1f}%** dari target — tidak ada gap yang perlu dianalisa lebih lanjut untuk periode/filter ini.")
else:
    gap_rp = tot_pendapatan_b - tot_pendapatan_r
    insights = []
    insights.append(
        f"Realisasi Pendapatan mencapai <b>{ach_pendapatan:.1f}%</b> dari target, dengan selisih (gap) sebesar <b>{fmt_rp(gap_rp)}</b>."
    )

    if ach_prestasi is not None:
        diff = ach_pendapatan - ach_prestasi
        if ach_prestasi < 100 and abs(diff) <= 10:
            insights.append(
                f"Prestasi juga hanya mencapai <b>{ach_prestasi:.1f}%</b> dari target — pola ini sejalan dengan capaian Pendapatan, "
                f"mengindikasikan bahwa <b>volume pekerjaan/prestasi unit yang belum tercapai</b> menjadi salah satu faktor utama rendahnya Pendapatan."
            )
        elif ach_prestasi < 100:
            lebih = "lebih kecil" if ach_prestasi > ach_pendapatan else "lebih besar"
            insights.append(
                f"Prestasi mencapai <b>{ach_prestasi:.1f}%</b> dari target — meski sama-sama di bawah target, gap-nya {lebih} dibanding Pendapatan, "
                f"sehingga faktor prestasi kemungkinan <b>{'turut berkontribusi namun bukan penyebab dominan' if ach_prestasi > ach_pendapatan else 'menjadi kontributor signifikan'}</b>."
            )
        else:
            insights.append(
                f"Prestasi justru mencapai <b>{ach_prestasi:.1f}%</b> (di atas target), sehingga rendahnya Pendapatan kemungkinan besar "
                f"<b>bukan disebabkan oleh volume pekerjaan/prestasi</b>, melainkan faktor lain seperti tarif/harga satuan atau piutang yang belum tertagih."
            )

    if target_populasi:
        gap_pop = target_populasi - realisasi_populasi
        gap_pop_pct = gap_pop / target_populasi * 100
        if gap_pop > 0:
            insights.append(
                f"Sebanyak <b>{gap_pop} unit ({gap_pop_pct:.1f}%)</b> dari total populasi yang ditargetkan <b>tidak mencatatkan realisasi Pendapatan sama sekali</b> "
                f"pada periode/filter ini — indikasi kuat adanya <b>unit yang tidak beroperasi (downtime/idle)</b>, yang turut menekan capaian Pendapatan secara keseluruhan."
            )
        else:
            insights.append(
                "Seluruh unit yang ditargetkan sudah mencatatkan realisasi Pendapatan (tidak ada indikasi unit idle/downtime dari sisi populasi)."
            )

    site_group = df.groupby("lokasi").agg(
        realisasi=("pendapatan_realisasi", "sum"),
        budget=("pendapatan_budget", "sum"),
    ).reset_index()
    site_group = site_group[site_group["budget"] > 0]
    if not site_group.empty:
        site_group["capaian"] = site_group["realisasi"] / site_group["budget"] * 100
        worst = site_group.sort_values("capaian").iloc[0]
        if worst["capaian"] < 100:
            insights.append(
                f"Site dengan capaian Pendapatan terendah adalah <b>{worst['lokasi']}</b> ({worst['capaian']:.1f}% dari target), "
                f"menjadi kontributor terbesar terhadap gap Pendapatan secara keseluruhan pada filter ini."
            )

    bullets_html = "".join([f"<li>{ins}</li>" for ins in insights])
    st.markdown(f'<div class="insight-box"><ul>{bullets_html}</ul></div>', unsafe_allow_html=True)

st.markdown("---")

# ---------------------------------------------------------------
# 7. ANALISA: UNIT TIDAK PRODUKTIF
# ---------------------------------------------------------------
st.markdown('<h3 class="section-title">Analisa: Unit Tidak Produktif</h3>', unsafe_allow_html=True)

# --- Agregasi per unit dari data utama: Pendapatan, Prestasi, Total Biaya ---
unit_fin = df.groupby(["id_unit", "nama_unit"], as_index=False).agg(
    lokasi=("lokasi", lambda s: s.mode().iat[0] if not s.mode().empty else s.iloc[0]),
    kategori=("kategori", lambda s: s.mode().iat[0] if not s.mode().empty else s.iloc[0]),
    pendapatan_r=("pendapatan_realisasi", "sum"),
    pendapatan_b=("pendapatan_budget", "sum"),
    prestasi_r=("prestasi_realisasi", "sum"),
    prestasi_b=("prestasi_budget", "sum"),
    total_biaya_r=("total_biaya_realisasi", "sum"),
)
unit_fin["unit_label"] = unit_fin["nama_unit"].apply(_unit_label)

# --- Agregasi maintenance per unit: frekuensi & biaya (scope site/bulan yang sama) ---
if not maint_raw.empty:
    maint_unit_agg = maint_df_site_bulan.copy()
    maint_unit_agg["unit_label"] = maint_unit_agg["nama_unit"].apply(_unit_label)
    maint_unit_agg = maint_unit_agg.groupby("unit_label", as_index=False).agg(
        maintenance_biaya=("biaya", "sum"),
        maintenance_freq=("biaya", "count"),
    )
else:
    maint_unit_agg = pd.DataFrame(columns=["unit_label", "maintenance_biaya", "maintenance_freq"])

unit_analysis = unit_fin.merge(maint_unit_agg, on="unit_label", how="left")
unit_analysis["maintenance_biaya"] = unit_analysis["maintenance_biaya"].fillna(0)
unit_analysis["maintenance_freq"] = unit_analysis["maintenance_freq"].fillna(0).astype(int)

unit_analysis["capaian_pendapatan"] = unit_analysis.apply(lambda r: achievement(r["pendapatan_r"], r["pendapatan_b"]), axis=1)
unit_analysis["capaian_prestasi"] = unit_analysis.apply(lambda r: achievement(r["prestasi_r"], r["prestasi_b"]), axis=1)
unit_analysis["margin"] = unit_analysis["pendapatan_r"] - unit_analysis["total_biaya_r"]

nonzero_freq = unit_analysis.loc[unit_analysis["maintenance_freq"] > 0, "maintenance_freq"]
if len(nonzero_freq) > 0:
    freq_p75 = nonzero_freq.quantile(0.75)
    unit_analysis["flag_maintenance_sering"] = unit_analysis["maintenance_freq"] >= freq_p75
else:
    unit_analysis["flag_maintenance_sering"] = False

unit_analysis["flag_pendapatan"] = unit_analysis["capaian_pendapatan"].apply(lambda v: (v is not None) and v < 100)
unit_analysis["flag_prestasi"] = unit_analysis["capaian_prestasi"].apply(lambda v: (v is not None) and v < 100)
unit_analysis["flag_margin_negatif"] = unit_analysis["margin"] < 0
unit_analysis["skor_masalah"] = (
    unit_analysis["flag_pendapatan"].astype(int)
    + unit_analysis["flag_prestasi"].astype(int)
    + unit_analysis["flag_maintenance_sering"].astype(int)
)
unit_analysis["tidak_produktif"] = unit_analysis["flag_margin_negatif"] & (unit_analysis["skor_masalah"] >= 2)

tp_df = unit_analysis[unit_analysis["tidak_produktif"]].sort_values("margin")

if tp_df.empty:
    st.success("Tidak ada unit yang teridentifikasi **Tidak Produktif** berdasarkan kriteria di atas pada filter saat ini.")
else:
    total_rugi = tp_df["margin"].sum()  # sudah negatif
    avg_capaian_pdt = tp_df["capaian_pendapatan"].mean()
    avg_freq_maint = tp_df["maintenance_freq"].mean()

    u1, u2, u3, u4 = st.columns(4)
    with u1:
        st.markdown(kpi_card(
            icon="⚠️", icon_bg=RED, accent=RED,
            label="Jumlah Unit Tidak Produktif",
            value=f"{len(tp_df):,} unit",
            budget_text=f"dari {len(unit_analysis):,} unit teranalisa",
            pill_text="Perlu Tindak Lanjut", pill_style="kpi-pill-red",
        ), unsafe_allow_html=True)
    with u2:
        st.markdown(kpi_card(
            icon="📉", icon_bg=RED, accent=RED,
            label="Total Margin Negatif (Kerugian)",
            value=fmt_rp(total_rugi),
            budget_text="Pendapatan − Total Biaya (unit tidak produktif)",
            pill_text="Rugi", pill_style="kpi-pill-red",
        ), unsafe_allow_html=True)
    with u3:
        st.markdown(kpi_card(
            icon="💰", icon_bg=GOLD, accent=GOLD,
            label="Rata-rata Capaian Pendapatan",
            value=(f"{avg_capaian_pdt:.1f}%" if pd.notna(avg_capaian_pdt) else "-"),
            budget_text="Rata-rata unit tidak produktif",
            pill_text="vs Budget", pill_style="kpi-pill-amber",
        ), unsafe_allow_html=True)
    with u4:
        st.markdown(kpi_card(
            icon="🔧", icon_bg=GOLD, accent=GOLD,
            label="Rata-rata Frekuensi Maintenance",
            value=f"{avg_freq_maint:.1f}x",
            budget_text="Rata-rata unit tidak produktif",
            pill_text="Sering", pill_style="kpi-pill-amber",
        ), unsafe_allow_html=True)

    st.markdown("##### Chart: Unit dengan Margin Ternegatif (Top 10)")
    top10 = tp_df.head(10).copy()
    top10["margin_jt"] = top10["margin"] / 1e6
    top10["label_margin"] = top10["margin"].apply(fmt_rp)
    top10["unit_chart_label"] = top10["id_unit"].astype(str) + " — " + top10["nama_unit"].astype(str)
    fig_tp = px.bar(
        top10.sort_values("margin_jt", ascending=False),
        y="unit_chart_label", x="margin_jt", orientation="h",
        text="label_margin",
        labels={"unit_chart_label": "Unit", "margin_jt": "Margin (Juta Rupiah)"},
    )
    fig_tp.update_traces(marker_color=RED, textposition="outside", textfont=dict(size=10, color=TEXT_LIGHT))
    fig_tp.update_layout(title="Top 10 Unit dengan Margin (Pendapatan − Biaya) Ternegatif", height=450, margin=dict(t=60, b=10, l=10))
    st.plotly_chart(style_fig(fig_tp), use_container_width=True)

    st.markdown("##### Rincian Unit Tidak Produktif")
    disp = tp_df.copy()
    disp["Status"] = disp["skor_masalah"].apply(
        lambda s: "🔴 Sangat Tidak Produktif (3 indikator)" if s == 3 else "🟠 Tidak Produktif (2 indikator)"
    )
    show_tp = pd.DataFrame({
        "ID Unit": disp["id_unit"],
        "Nama Unit": disp["nama_unit"],
        "Site": disp["lokasi"],
        "Kategori": disp["kategori"].map(lambda k: KATEGORI_LABEL.get(k, k)),
        "Pendapatan Realisasi": disp["pendapatan_r"].apply(fmt_rp),
        "Capaian Pendapatan": disp["capaian_pendapatan"].apply(lambda v: f"{v:.1f}%" if pd.notna(v) else "-"),
        "Capaian Prestasi": disp["capaian_prestasi"].apply(lambda v: f"{v:.1f}%" if pd.notna(v) else "-"),
        "Total Biaya": disp["total_biaya_r"].apply(fmt_rp),
        "Margin (Pendapatan − Biaya)": disp["margin"].apply(fmt_rp),
        "Frekuensi Maintenance": disp["maintenance_freq"],
        "Biaya Maintenance": disp["maintenance_biaya"].apply(fmt_rp),
        "Status": disp["Status"],
    })
    st.dataframe(show_tp, use_container_width=True, hide_index=True, height=420)
    csv_tp = show_tp.to_csv(index=False).encode("utf-8")
    st.download_button("⬇️ Unduh Daftar Unit Tidak Produktif (CSV)", csv_tp, file_name="unit_tidak_produktif_bkms.csv", mime="text/csv")

    worst3 = tp_df.head(3)
    poin = []
    for _, r in worst3.iterrows():
        cap_pdt = f"{r['capaian_pendapatan']:.1f}%" if pd.notna(r["capaian_pendapatan"]) else "tanpa budget"
        cap_prs = f"{r['capaian_prestasi']:.1f}%" if pd.notna(r["capaian_prestasi"]) else "tanpa budget"
        poin.append(
            f"<b>{r['id_unit']} — {r['nama_unit']}</b> (site {r['lokasi']}): margin <b>{fmt_rp(r['margin'])}</b>, "
            f"capaian Pendapatan {cap_pdt}, capaian Prestasi {cap_prs}, maintenance {int(r['maintenance_freq'])} kali "
            f"({fmt_rp(r['maintenance_biaya'])})."
        )
    bullets_tp_html = "".join([f"<li>{p}</li>" for p in poin])
    st.markdown(f'<div class="insight-box"><b>3 unit paling bermasalah:</b><ul>{bullets_tp_html}</ul></div>', unsafe_allow_html=True)

st.markdown("---")
st.caption("Dashboard Operational Review • PT Buana Karya Mandiri Sejahtera (BKMS) • Dibuat oleh ALIP BA TA")
